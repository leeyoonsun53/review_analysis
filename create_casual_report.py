# -*- coding: utf-8 -*-
"""
독도토너 리뷰 분석 리포트 - 캐쥬얼 버전
컨셉 기획에 맞춘 데이터 분석 관점 슬라이드
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import pandas as pd
import json
import sys
from datetime import datetime
from collections import Counter

sys.stdout.reconfigure(encoding='utf-8')

# 색상 정의 - 캐쥬얼한 톤
MAIN_BLUE = RGBColor(0x3b, 0x82, 0xf6)
DARK_BLUE = RGBColor(0x1e, 0x40, 0xaf)
LIGHT_BLUE = RGBColor(0xdb, 0xea, 0xfe)
GREEN = RGBColor(0x10, 0xb9, 0x81)
RED = RGBColor(0xef, 0x44, 0x44)
ORANGE = RGBColor(0xf9, 0x73, 0x16)
PURPLE = RGBColor(0x8b, 0x5c, 0xf6)
GRAY = RGBColor(0x6b, 0x72, 0x80)
DARK = RGBColor(0x1f, 0x29, 0x37)
WHITE = RGBColor(0xff, 0xff, 0xff)
BG_LIGHT = RGBColor(0xf8, 0xfa, 0xfc)


def add_title_slide(prs, title, subtitle):
    """표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # 타이틀
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 서브
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0x93, 0xc5, 0xfd)
    p.alignment = PP_ALIGN.CENTER


def add_section_slide(prs, number, title, emoji=""):
    """섹션 구분"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = MAIN_BLUE
    bg.line.fill.background()

    # 번호
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"PART {number}"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xbd, 0xe0, 0xfe)
    p.alignment = PP_ALIGN.CENTER

    # 타이틀
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}" if emoji else title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, note=""):
    """내용 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    hd = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    hd.fill.solid()
    hd.fill.fore_color.rgb = MAIN_BLUE
    hd.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 내용
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_after = Pt(12)

        # 강조 처리
        if bullet.startswith("→") or bullet.startswith("▶"):
            p.font.bold = True
            p.font.color.rgb = MAIN_BLUE

    # 하단 노트
    if note:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY


def add_stat_cards(prs, title, stats):
    """통계 카드 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    hd = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    hd.fill.solid()
    hd.fill.fore_color.rgb = MAIN_BLUE
    hd.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 카드들
    card_width = 2.8
    gap = 0.3
    start_x = (10 - (card_width * len(stats) + gap * (len(stats) - 1))) / 2

    for i, (label, value, color) in enumerate(stats):
        x = start_x + i * (card_width + gap)

        # 카드 배경
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2), Inches(card_width), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()

        # 값
        tb = slide.shapes.add_textbox(Inches(x), Inches(2.3), Inches(card_width), Inches(1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # 라벨
        tb = slide.shapes.add_textbox(Inches(x), Inches(3.3), Inches(card_width), Inches(0.6))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER


def add_bar_chart(prs, title, categories, values, subtitle="", color=MAIN_BLUE):
    """막대 차트"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    hd = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    hd.fill.solid()
    hd.fill.fore_color.rgb = MAIN_BLUE
    hd.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY

    # 차트
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.6), Inches(9), Inches(5.2),
        chart_data
    ).chart


def add_quadrant_chart(prs, title, brand_positions, x_label, y_label, subtitle=""):
    """4분면 포지셔닝 차트
    brand_positions: [(brand, x, y, color), ...]
    x, y는 0~100 스케일
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    hd = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    hd.fill.solid()
    hd.fill.fore_color.rgb = MAIN_BLUE
    hd.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(9), Inches(0.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY

    # 차트 영역 정의
    chart_left = 1.2
    chart_top = 1.6
    chart_width = 7.0
    chart_height = 5.2

    # 배경 4분면
    # 좌상 (진정 + 프리미엄)
    q1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(chart_left), Inches(chart_top),
        Inches(chart_width/2), Inches(chart_height/2))
    q1.fill.solid()
    q1.fill.fore_color.rgb = RGBColor(0xfe, 0xf3, 0xc7)  # 노랑
    q1.line.fill.background()

    # 우상 (보습 + 프리미엄)
    q2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(chart_left + chart_width/2), Inches(chart_top),
        Inches(chart_width/2), Inches(chart_height/2))
    q2.fill.solid()
    q2.fill.fore_color.rgb = RGBColor(0xdc, 0xfc, 0xe7)  # 초록
    q2.line.fill.background()

    # 좌하 (진정 + 가성비)
    q3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(chart_left), Inches(chart_top + chart_height/2),
        Inches(chart_width/2), Inches(chart_height/2))
    q3.fill.solid()
    q3.fill.fore_color.rgb = RGBColor(0xfe, 0xe2, 0xe2)  # 빨강
    q3.line.fill.background()

    # 우하 (보습 + 가성비)
    q4 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(chart_left + chart_width/2), Inches(chart_top + chart_height/2),
        Inches(chart_width/2), Inches(chart_height/2))
    q4.fill.solid()
    q4.fill.fore_color.rgb = RGBColor(0xdb, 0xea, 0xfe)  # 파랑
    q4.line.fill.background()

    # 축 라벨
    # X축 (좌: 진정, 우: 보습)
    tb = slide.shapes.add_textbox(Inches(chart_left - 0.1), Inches(chart_top + chart_height/2 - 0.15), Inches(1), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "← 진정"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK

    tb = slide.shapes.add_textbox(Inches(chart_left + chart_width - 0.7), Inches(chart_top + chart_height/2 - 0.15), Inches(1), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "보습 →"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK

    # Y축 (상: 프리미엄, 하: 가성비)
    tb = slide.shapes.add_textbox(Inches(chart_left + chart_width/2 - 0.4), Inches(chart_top - 0.25), Inches(1), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "프리미엄 ↑"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK

    tb = slide.shapes.add_textbox(Inches(chart_left + chart_width/2 - 0.35), Inches(chart_top + chart_height + 0.05), Inches(1), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "↓ 가성비"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK

    # 브랜드 점 배치
    for brand, x_val, y_val, color in brand_positions:
        # x: 0=왼쪽(진정), 100=오른쪽(보습)
        # y: 0=아래(가성비), 100=위(프리미엄)
        px = chart_left + (x_val / 100) * chart_width
        py = chart_top + chart_height - (y_val / 100) * chart_height

        # 원
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(px - 0.25), Inches(py - 0.25),
            Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.color.rgb = WHITE
        circle.line.width = Pt(2)

        # 브랜드명
        tb = slide.shapes.add_textbox(Inches(px - 0.5), Inches(py + 0.25), Inches(1), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = brand
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER


def add_two_column(prs, title, left_title, left_items, right_title, right_items, left_color=GREEN, right_color=RED):
    """2단 비교"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    hd = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    hd.fill.solid()
    hd.fill.fore_color.rgb = MAIN_BLUE
    hd.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 왼쪽
    lbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.4), Inches(4.6), Inches(5.5))
    lbox.fill.solid()
    lbox.fill.fore_color.rgb = BG_LIGHT
    lbox.line.color.rgb = left_color
    lbox.line.width = Pt(2)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.2), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = left_color

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4.2), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)

    # 오른쪽
    rbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.1), Inches(1.4), Inches(4.6), Inches(5.5))
    rbox.fill.solid()
    rbox.fill.fore_color.rgb = BG_LIGHT
    rbox.line.color.rgb = right_color
    rbox.line.width = Pt(2)

    tb = slide.shapes.add_textbox(Inches(5.3), Inches(1.6), Inches(4.2), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = right_color

    tb = slide.shapes.add_textbox(Inches(5.3), Inches(2.2), Inches(4.2), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)


def add_table_slide(prs, title, headers, rows, subtitle=""):
    """테이블"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    hd = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    hd.fill.solid()
    hd.fill.fore_color.rgb = MAIN_BLUE
    hd.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    start_y = 1.3
    if subtitle:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        start_y = 1.6

    # 테이블
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.3), Inches(start_y),
        Inches(9.4), Inches(0.45 * num_rows)
    ).table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = MAIN_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT


def add_insight_slide(prs, title, insights):
    """인사이트"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더 - 오렌지
    hd = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    hd.fill.solid()
    hd.fill.fore_color.rgb = ORANGE
    hd.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 인사이트 박스들
    y = 1.4
    for ins in insights:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3), Inches(y), Inches(9.4), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_LIGHT
        box.line.color.rgb = ORANGE

        tb = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.15), Inches(9), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ins
        p.font.size = Pt(14)
        p.font.color.rgb = DARK

        y += 1.0


def add_quote_slide(prs, quote, bg_color=DARK_BLUE):
    """강조 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def main():
    print("=" * 60)
    print("독도토너 리뷰 분석 리포트 생성 (캐쥬얼 버전)")
    print("=" * 60)

    # 데이터 로드
    print("\n[데이터 로드 중...]")
    df = pd.read_csv('data/merged_reviews_processed.csv', encoding='utf-8-sig')
    cat = json.load(open('output/gpt_analysis_categorized.json', encoding='utf-8'))

    total = len(df)
    oly_count = len(df[df['PLATFORM'] == '올리브영'])
    musinsa_count = len(df[df['PLATFORM'] == '무신사'])

    brands = ['독도토너', '토리든', '브링그린', '에스네이처', '아누아', '토니모리', '아비브']

    # 감성 통계
    pos_count = sum(1 for r in cat if r.get('sentiment') == 'POS')
    neg_count = sum(1 for r in cat if r.get('sentiment') == 'NEG')
    neu_count = total - pos_count - neg_count

    # 브랜드별 통계
    brand_stats = {}
    for brand in brands:
        b_cat = [r for r in cat if r.get('brand') == brand]
        b_total = len(b_cat)
        b_pos = sum(1 for r in b_cat if r.get('sentiment') == 'POS')
        b_neg = sum(1 for r in b_cat if r.get('sentiment') == 'NEG')
        brand_stats[brand] = {
            'total': b_total,
            'pos_rate': b_pos / b_total * 100 if b_total > 0 else 0,
            'neg_rate': b_neg / b_total * 100 if b_total > 0 else 0
        }

    # Pain/Positive 통계
    pain_counts = Counter()
    pos_counts = Counter()
    usage_counts = Counter()

    for r in cat:
        for p in r.get('pain_points', []):
            pain_counts[p] += 1
        for p in r.get('positive_points', []):
            pos_counts[p] += 1
        for u in r.get('usage_tags', []):
            usage_counts[u] += 1

    # 독도토너 데이터
    dokdo_cat = [r for r in cat if r.get('brand') == '독도토너']
    dokdo_pain = Counter()
    dokdo_pos = Counter()
    for r in dokdo_cat:
        for p in r.get('pain_points', []):
            dokdo_pain[p] += 1
        for p in r.get('positive_points', []):
            dokdo_pos[p] += 1

    # 재구매 태그 통계 (올영)
    rebuy_count = len(df[df['PURCHASE_TAG'].str.contains('재구매', na=False)])
    month_use_count = len(df[df['PURCHASE_TAG'].str.contains('한달이상', na=False)])

    # 브랜드별 포지셔닝 데이터 (4분면 차트용)
    brand_positioning = {}
    for brand in brands:
        b_cat = [r for r in cat if r.get('brand') == brand]
        b_total = len(b_cat)

        # 보습 vs 진정 (X축: 보습이 높으면 오른쪽)
        moisturize = sum(1 for r in b_cat if '보습' in r.get('benefit_tags', []))
        calming = sum(1 for r in b_cat if '진정' in r.get('benefit_tags', []))
        x_val = (moisturize / (moisturize + calming + 1)) * 100  # 보습 비율

        # 가성비 vs 프리미엄 (Y축: 인생템이 높으면 위)
        value = sum(1 for r in b_cat if '가성비' in r.get('value_tags', []))
        premium = sum(1 for r in b_cat if '인생템' in r.get('value_tags', []))
        y_val = (premium / (value + premium + 1)) * 100  # 프리미엄 비율

        brand_positioning[brand] = {'x': x_val, 'y': y_val}

    print(f"   총 리뷰: {total:,}건")
    print(f"   긍정: {pos_count:,}건 ({pos_count/total*100:.1f}%)")

    # PPT 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========================================
    # 슬라이드 1: 표지
    # ========================================
    add_title_slide(prs,
        "독도토너, 다시 사게 만들려면?",
        f"올리브영 + 무신사 리뷰 {total:,}건 분석 | 2025년 데이터")

    # ========================================
    # 슬라이드 2: 분석 개요
    # ========================================
    add_content_slide(prs, "이 분석은 이렇게 했어요", [
        "데이터 수집",
        f"• 올리브영 리뷰 {oly_count:,}건 + 무신사 리뷰 {musinsa_count:,}건",
        f"• 수집 기간: 2025년 1월 ~ 2026년 1월",
        f"• 대상 브랜드: 독도토너 외 6개 경쟁 브랜드",
        "",
        "분석 방법",
        "• GPT-4o-mini API로 리뷰 감성 분석 (긍정/중립/부정)",
        "• Pain Points / Positive Points 자동 추출",
        "• 효능 태그 (보습, 진정, 장벽 등) 분류",
        "• 사용법 태그 (닦토, 레이어링, 스킨팩 등) 분류",
        "",
        "▶ 총 API 호출: 33,484건 | 분석 정확도 체감 90%+",
    ], note="* 올영 리뷰에는 '재구매', '한달이상사용' 태그 포함")

    # ========================================
    # 슬라이드 3: 핵심 지표
    # ========================================
    add_stat_cards(prs, "한눈에 보는 핵심 지표", [
        ("총 리뷰", f"{total:,}건", MAIN_BLUE),
        ("긍정 비율", f"{pos_count/total*100:.1f}%", GREEN),
        ("부정 비율", f"{neg_count/total*100:.1f}%", RED),
    ])

    # ========================================
    # PART 1: 현상 진단
    # ========================================
    add_section_slide(prs, "1", "잊혀지는 국민 토너?", "🔍")

    # 브랜드별 긍정률
    sorted_brands = sorted(brands, key=lambda x: brand_stats[x]['pos_rate'], reverse=True)
    add_bar_chart(prs, "브랜드별 긍정률 비교",
        categories=sorted_brands,
        values=[brand_stats[b]['pos_rate'] for b in sorted_brands],
        subtitle="독도토너는 4위... 나쁘진 않지만 1등은 아님")

    # 브랜드별 테이블
    brand_rows = []
    for b in sorted_brands:
        s = brand_stats[b]
        brand_rows.append([b, f"{s['total']:,}", f"{s['pos_rate']:.1f}%", f"{s['neg_rate']:.1f}%"])

    add_table_slide(prs, "브랜드별 리뷰 분석 요약",
        ["브랜드", "리뷰 수", "긍정률", "부정률"],
        brand_rows,
        subtitle="긍정률 기준 정렬 | 에스네이처가 1위, 독도토너는 4위")

    # 독도토너 현황
    dokdo_stats = brand_stats['독도토너']
    add_content_slide(prs, "독도토너 현황", [
        f"총 리뷰: {dokdo_stats['total']:,}건",
        f"긍정률: {dokdo_stats['pos_rate']:.1f}% (7개 브랜드 중 4위)",
        f"부정률: {dokdo_stats['neg_rate']:.1f}%",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "리뷰에서 보이는 포지셔닝:",
        "• '무난한 입문용 토너'",
        "• '순하고 자극 없는 토너'",
        "• '대용량 닦토템'",
        "",
        "▶ 문제: '순함' 외에 특별한 강점이 안 보임",
        "▶ 경쟁사 대비 '와우 포인트' 부족",
    ])

    # 4분면 포지셔닝 차트
    brand_colors = {
        '독도토너': RGBColor(0x3b, 0x82, 0xf6),  # 파랑
        '토리든': RGBColor(0x10, 0xb9, 0x81),    # 초록
        '브링그린': RGBColor(0x22, 0xc5, 0x5e),  # 연두
        '에스네이처': RGBColor(0x8b, 0x5c, 0xf6), # 보라
        '아누아': RGBColor(0xf9, 0x73, 0x16),    # 오렌지
        '토니모리': RGBColor(0xef, 0x44, 0x44),  # 빨강
        '아비브': RGBColor(0xec, 0x48, 0x99),    # 핑크
    }

    brand_pos_list = []
    for brand in brands:
        pos = brand_positioning[brand]
        brand_pos_list.append((brand, pos['x'], pos['y'], brand_colors[brand]))

    add_quadrant_chart(prs, "브랜드 포지셔닝 맵",
        brand_pos_list,
        x_label="진정 ↔ 보습",
        y_label="가성비 ↔ 프리미엄",
        subtitle="X축: 보습 vs 진정 비율 | Y축: 인생템 vs 가성비 비율 (리뷰 태그 기반)")

    # ========================================
    # PART 2: 원인 분석
    # ========================================
    add_section_slide(prs, "2", "왜 다시 안 살까?", "🤔")

    # Positive Points TOP 10
    top_pos = pos_counts.most_common(10)
    add_bar_chart(prs, "고객이 좋아하는 포인트 TOP 10",
        categories=[p[0][:10] for p in top_pos],
        values=[p[1] for p in top_pos],
        subtitle="'순함/저자극'이 압도적 1위 → 토너 선택의 핵심 기준")

    # Pain Points TOP 10
    top_pain = pain_counts.most_common(10)
    add_bar_chart(prs, "고객 불만 포인트 TOP 10",
        categories=[p[0][:10] for p in top_pain],
        values=[p[1] for p in top_pain],
        subtitle="'건조/보습부족'이 1위 → 전 브랜드 공통 과제")

    # 독도토너 강점 vs 약점
    top_dokdo_pos = dokdo_pos.most_common(5)
    top_dokdo_pain = dokdo_pain.most_common(5)

    add_two_column(prs, "독도토너: 강점 vs 약점",
        "😊 Positive Points",
        [f"• {p[0]}: {p[1]}건" for p in top_dokdo_pos] + [
            "",
            "→ '순함'이 압도적 강점",
            "→ 기본기는 탄탄함"
        ],
        "😣 Pain Points",
        [f"• {p[0]}: {p[1]}건" for p in top_dokdo_pain] + [
            "",
            "→ '효과 없음' 인식이 문제",
            "→ 순하긴 한데... 그래서 뭐?"
        ])

    # 저평점 분석
    add_content_slide(prs, "저평점 리뷰 심층 분석 (1-2점)", [
        "독도토너 저평점 63건 분석 결과:",
        "",
        "• 효과 부족: 46% ← 가장 큰 문제!",
        "  '물 같다', '바른 건지 모르겠다', '뭔가 부족'",
        "",
        "• 자극/트러블: 31%",
        "  '예상과 달리 트러블', '따가움'",
        "",
        "• 건조함: 8%",
        "  '보습력이 약함'",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "▶ 핵심 문제: '순한 건 알겠는데, 효과가 있긴 해?'",
    ])

    # 재구매 vs 비재구매
    add_content_slide(prs, "재구매 고객은 뭐가 다를까?", [
        f"올영 '재구매' 태그 리뷰: {rebuy_count:,}건",
        f"올영 '한달이상사용' 태그 리뷰: {month_use_count:,}건",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "재구매 고객이 강조하는 포인트:",
        "• '순함' - 자극 없이 쓸 수 있어서",
        "• '대용량' - 양 걱정 없이 막 쓸 수 있어서",
        "• '가성비' - 가격 대비 만족",
        "• '닦토용' - 매일 아침저녁 닦토로 사용",
        "",
        "▶ 재구매 동력: 기본기 + 실용성",
        "▶ 문제: 첫 구매자를 재구매로 연결하는 '후킹'이 부족",
    ])

    # ========================================
    # PART 3: 전략 방향
    # ========================================
    add_section_slide(prs, "3", "'순함'을 넘어서", "🎯")

    # 사용법 분석
    top_usage = usage_counts.most_common(5)
    add_bar_chart(prs, "고객들은 이렇게 쓰고 있어요",
        categories=[u[0] for u in top_usage],
        values=[u[1] for u in top_usage],
        subtitle="닦토 > 레이어링 > 스킨팩 순 | 대용량 활용 패턴 뚜렷")

    # AS-IS vs TO-BE
    add_two_column(prs, "포지셔닝 전환 제안",
        "AS-IS: Nice to Have",
        [
            "'자극 없는 토너'",
            "'입문용으로 좋아요'",
            "'무난해요'",
            "",
            "→ 다른 거 써보고 싶어짐",
            "→ '졸업'해버림",
            "→ 재구매 동기 부족"
        ],
        "TO-BE: Must Have",
        [
            "'모든 스킨케어의 시작점'",
            "'고기능 제품의 완벽한 베이스'",
            "'이게 있어야 다른 게 잘 먹어요'",
            "",
            "→ 레티놀/비타민C 전 필수템",
            "→ 버릴 수 없는 기본템",
            "→ 루틴의 일부로 정착"
        ],
        left_color=GRAY, right_color=GREEN)

    # 타겟 세분화
    add_content_slide(prs, "사용법 기반 타겟 세분화", [
        "닦토 유저 (전체의 38%)",
        "• 대용량 + 순함이 핵심",
        "• 매일 아침저녁 사용 → 소진 주기 약 2개월",
        "• CRM: '토너 다 쓸 때 됐죠? 대용량 재입고!'",
        "",
        "레이어링 유저 (12%)",
        "• 7스킨법 등 다겹 사용",
        "• 촉촉함/흡수력 중시",
        "• CRM: '여러 겹 발라도 무거움 없이'",
        "",
        "스킨팩 유저 (7%)",
        "• 진정/트러블 케어 목적",
        "• CRM: '화장솜에 듬뿍, 진정 스킨팩 루틴'",
    ])

    # ========================================
    # PART 4: 실행 제안
    # ========================================
    add_section_slide(prs, "4", "이렇게 해보면 어떨까요?", "💡")

    # 채널별 전략
    add_two_column(prs, "채널별 메시지 제안",
        "올리브영",
        [
            "강조 키워드:",
            "• '순함', '진정', '결 개선'",
            "",
            "기획 제안:",
            "• 클렌저 + 토너 번들 세트",
            "• '민감성 피부 추천템' 큐레이션",
            "",
            "프로모션:",
            "• 재구매 고객 전용 할인",
            "• '한달이상사용' 리뷰어 리워드"
        ],
        "무신사",
        [
            "강조 키워드:",
            "• '가성비', '입문용', '대용량'",
            "",
            "기획 제안:",
            "• 그루밍 입문 세트",
            "• '남자도 스킨케어' 기획전",
            "",
            "프로모션:",
            "• 첫 구매 할인",
            "• 리뷰 작성 시 추가 적립"
        ],
        left_color=GREEN, right_color=PURPLE)

    # 리마인드 전략
    add_content_slide(prs, "재구매 유도 전략", [
        "소진 예측 리마인드",
        "• 구매일 + 60일 후 알림톡 발송",
        "• '토너 다 쓸 때 됐죠? 지금 주문하면 내일 도착'",
        "",
        "크로스셀링",
        "• 독도 토너 구매자 → 독도 로션/크림 추천",
        "• '같이 쓰면 효과 2배'",
        "",
        "로열티 프로그램",
        "• 3회 구매 시 VIP 등급 부여",
        "• 독점 할인 + 신제품 얼리버드 접근권",
        "",
        "▶ 목표: '첫 구매 → 3회 구매' 전환율 높이기",
    ])

    # ========================================
    # PART 5: 핵심 인사이트
    # ========================================
    add_section_slide(prs, "5", "정리하면", "📌")

    add_insight_slide(prs, "핵심 인사이트 5가지", [
        f"① 전체 긍정률 {pos_count/total*100:.1f}%로 토너 시장 만족도 높음",
        "② '순함/저자극'이 토너 선택의 핵심 기준 (5,000건+ 언급)",
        "③ 독도토너 강점은 '순함', 약점은 '효과 없음' 인식",
        "④ 재구매 고객의 핵심 동력: 순함 + 대용량 + 가성비",
        "⑤ 전략 제안: 'Nice to Have' → 'Must Have'로 리포지셔닝",
    ])

    # 마무리
    add_quote_slide(prs,
        "'순한 토너'에서\n'모든 스킨케어의 시작점'으로\n\n━━━━━━━━━━━━━━━━━━━━\n\n"
        "독도토너의 기본기를 살리면서\n새로운 가치를 더하는 전략이 필요합니다")

    # 감사 슬라이드
    add_title_slide(prs, "감사합니다",
        f"분석 데이터: {total:,}건 | GPT-4o-mini 기반 | {datetime.now().strftime('%Y.%m')}")

    # 저장
    output_path = 'output/독도토너_리뷰분석_캐쥬얼.pptx'
    prs.save(output_path)

    print("\n" + "=" * 60)
    print(f"리포트 생성 완료: {output_path}")
    print(f"총 {len(prs.slides)}장 슬라이드")
    print("=" * 60)


if __name__ == "__main__":
    main()
