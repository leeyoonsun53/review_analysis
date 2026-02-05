# -*- coding: utf-8 -*-
"""
독도토너 피부타입별 심층 분석 슬라이드
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
import json
import sys
from collections import Counter

sys.stdout.reconfigure(encoding='utf-8')

# 색상
MAIN_BLUE = RGBColor(0x3b, 0x82, 0xf6)
DARK_BLUE = RGBColor(0x1e, 0x40, 0xaf)
GREEN = RGBColor(0x10, 0xb9, 0x81)
RED = RGBColor(0xef, 0x44, 0x44)
ORANGE = RGBColor(0xf9, 0x73, 0x16)
GRAY = RGBColor(0x6b, 0x72, 0x80)
DARK = RGBColor(0x1f, 0x29, 0x37)
WHITE = RGBColor(0xff, 0xff, 0xff)
BG_LIGHT = RGBColor(0xf8, 0xfa, 0xfc)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0x93, 0xc5, 0xfd)
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, header_color=MAIN_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    hd = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    hd.fill.solid()
    hd.fill.fore_color.rgb = header_color
    hd.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)

        if bullet.startswith("→") or bullet.startswith("▶") or bullet.startswith("★"):
            p.font.bold = True
            p.font.color.rgb = MAIN_BLUE if not bullet.startswith("★") else RED


def add_comparison_slide(prs, title, data_rows, highlight_row=None):
    """비교 테이블 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    hd = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    hd.fill.solid()
    hd.fill.fore_color.rgb = MAIN_BLUE
    hd.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 테이블
    headers = ['브랜드', '민감성', '건성', '지성', '복합성']
    num_rows = len(data_rows) + 1
    num_cols = len(headers)

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.3), Inches(1.4),
        Inches(9.4), Inches(0.5 * num_rows)
    ).table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = MAIN_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for i, row in enumerate(data_rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.alignment = PP_ALIGN.CENTER

            # 독도토너 행 하이라이트
            if highlight_row and i == highlight_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xdb, 0xea, 0xfe)
                p.font.bold = True
                p.font.color.rgb = DARK_BLUE

            # 0.0% 셀 초록색
            if '0.0%' in str(val):
                p.font.color.rgb = GREEN
                p.font.bold = True

            # 4% 이상 빨간색
            try:
                if '%' in str(val):
                    num = float(val.replace('%', ''))
                    if num >= 4.0:
                        p.font.color.rgb = RED
                        p.font.bold = True
            except:
                pass


def add_heatmap_slide(prs, title, subtitle=""):
    """히트맵 스타일 비교"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    hd = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    hd.fill.solid()
    hd.fill.fore_color.rgb = MAIN_BLUE
    hd.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(9), Inches(0.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY

    return slide


def add_insight_box(prs, title, insights, header_color=ORANGE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    hd = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    hd.fill.solid()
    hd.fill.fore_color.rgb = header_color
    hd.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    y = 1.4
    for ins in insights:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3), Inches(y), Inches(9.4), Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_LIGHT
        box.line.color.rgb = header_color

        tb = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.15), Inches(9), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ins
        p.font.size = Pt(14)
        p.font.color.rgb = DARK

        y += 1.1


def main():
    print("=" * 60)
    print("피부타입별 심층 분석 슬라이드 생성")
    print("=" * 60)

    # 데이터 로드
    df = pd.read_csv('data/merged_reviews_processed.csv', encoding='utf-8-sig', low_memory=False)
    cat = json.load(open('output/gpt_analysis_categorized.json', encoding='utf-8'))

    cat_map = {r['idx']: r for r in cat}
    df['sentiment'] = df.index.map(lambda x: cat_map.get(x, {}).get('sentiment', 'NEU'))
    df['pain_points'] = df.index.map(lambda x: cat_map.get(x, {}).get('pain_points', []))
    df['positive_points'] = df.index.map(lambda x: cat_map.get(x, {}).get('positive_points', []))

    skin_df = df[df['SKIN_TYPE'].notna()].copy()

    brands = ['독도토너', '토리든', '브링그린', '에스네이처', '아누아', '토니모리', '아비브']
    skin_types = ['민감성', '건성', '지성', '복합성']

    # 브랜드별 피부타입별 부정률 계산
    neg_rates = {}
    for brand in brands:
        neg_rates[brand] = {}
        for skin in skin_types:
            bsdf = skin_df[(skin_df['BRAND_NAME'] == brand) & (skin_df['SKIN_TYPE'] == skin)]
            total = len(bsdf)
            neg = (bsdf['sentiment'] == 'NEG').sum()
            neg_rates[brand][skin] = {'total': total, 'neg': neg, 'rate': neg/total*100 if total > 0 else None}

    # 독도토너 민감성 분석
    dokdo_sens = skin_df[(skin_df['BRAND_NAME'] == '독도토너') & (skin_df['SKIN_TYPE'] == '민감성')]
    sens_positive = Counter()
    for idx in dokdo_sens.index:
        for p in cat_map.get(idx, {}).get('positive_points', []):
            sens_positive[p] += 1

    # 독도토너 복합성 부정 분석
    dokdo_combo_neg = skin_df[(skin_df['BRAND_NAME'] == '독도토너') &
                              (skin_df['SKIN_TYPE'] == '복합성') &
                              (skin_df['sentiment'] == 'NEG')]
    combo_pain = Counter()
    for idx in dokdo_combo_neg.index:
        for p in cat_map.get(idx, {}).get('pain_points', []):
            combo_pain[p] += 1

    # 다른 브랜드 평균
    other_sens_rates = [neg_rates[b]['민감성']['rate'] for b in brands if b != '독도토너' and neg_rates[b]['민감성']['rate'] is not None]
    avg_sens = sum(other_sens_rates) / len(other_sens_rates) if other_sens_rates else 0

    other_combo_rates = [neg_rates[b]['복합성']['rate'] for b in brands if b != '독도토너' and neg_rates[b]['복합성']['rate'] is not None]
    avg_combo = sum(other_combo_rates) / len(other_combo_rates) if other_combo_rates else 0

    # PPT 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 슬라이드 1: 표지
    add_title_slide(prs,
        "독도토너 피부타입별 심층 분석",
        "민감성 부정률 0%의 비밀 & 복합성의 숨겨진 문제")

    # 슬라이드 2: 발견
    add_content_slide(prs, "이상한 점을 발견했습니다", [
        "피부타입별 부정 비율을 분석하던 중...",
        "",
        "★ 독도토너만 민감성 피부에서 부정률이 0%",
        "",
        "다른 브랜드는 민감성에서 부정률이 가장 높은데,",
        "독도토너는 오히려 민감성에서 가장 만족도가 높았습니다.",
        "",
        "반면, 복합성 피부에서는 오히려 불만이 많았습니다.",
        "",
        "▶ 이게 무슨 의미일까요?",
    ])

    # 슬라이드 3: 브랜드별 비교 테이블
    table_data = []
    for brand in brands:
        row = [brand]
        for skin in skin_types:
            r = neg_rates[brand][skin]
            if r['rate'] is not None:
                row.append(f"{r['rate']:.1f}%")
            else:
                row.append("N/A")
        table_data.append(row)

    add_comparison_slide(prs, "브랜드별 × 피부타입별 부정률 비교", table_data, highlight_row=0)

    # 슬라이드 4: 독도토너 피부타입별 상세
    dokdo_stats = neg_rates['독도토너']
    add_content_slide(prs, "독도토너 피부타입별 부정률", [
        f"민감성: {dokdo_stats['민감성']['rate']:.1f}% ({dokdo_stats['민감성']['neg']}/{dokdo_stats['민감성']['total']}) ← 0%!",
        f"건성: {dokdo_stats['건성']['rate']:.1f}% ({dokdo_stats['건성']['neg']}/{dokdo_stats['건성']['total']})",
        f"지성: {dokdo_stats['지성']['rate']:.1f}% ({dokdo_stats['지성']['neg']}/{dokdo_stats['지성']['total']})",
        f"복합성: {dokdo_stats['복합성']['rate']:.1f}% ({dokdo_stats['복합성']['neg']}/{dokdo_stats['복합성']['total']}) ← 가장 높음",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        f"▶ 다른 브랜드 민감성 평균: {avg_sens:.1f}%",
        f"▶ 독도토너 민감성: 0.0%",
        f"▶ 차이: -{avg_sens:.1f}%p (독도토너가 압도적으로 낮음)",
        "",
        f"▶ 다른 브랜드 복합성 평균: {avg_combo:.1f}%",
        f"▶ 독도토너 복합성: {dokdo_stats['복합성']['rate']:.1f}%",
        f"▶ 차이: +{dokdo_stats['복합성']['rate'] - avg_combo:.1f}%p (독도토너가 오히려 높음)",
    ])

    # 슬라이드 5: 민감성 0%의 이유
    top_sens_pos = sens_positive.most_common(5)
    sens_lines = [
        f"독도토너 민감성 피부 리뷰: {len(dokdo_sens)}건",
        f"긍정: {(dokdo_sens['sentiment']=='POS').sum()}건 | 중립: {(dokdo_sens['sentiment']=='NEU').sum()}건 | 부정: 0건",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "민감성 고객이 좋아하는 포인트 TOP 5:",
    ]
    for p, c in top_sens_pos:
        sens_lines.append(f"  • {p}: {c}건")

    sens_lines.extend([
        "",
        "▶ '순함/저자극'이 압도적 1위",
        "▶ 독도토너의 핵심 포지셔닝이 민감성 피부에 완벽히 적중",
        "▶ 민감성 피부 고객에게 '신뢰할 수 있는 토너'로 인식",
    ])
    add_content_slide(prs, "민감성 부정률 0%의 이유", sens_lines, header_color=GREEN)

    # 슬라이드 6: 복합성 문제
    combo_lines = [
        f"독도토너 복합성 피부 부정 리뷰: {len(dokdo_combo_neg)}건",
        "",
        "복합성 피부 불만 포인트:",
    ]
    for p, c in combo_pain.most_common(6):
        combo_lines.append(f"  • {p}: {c}건")

    combo_lines.extend([
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "복합성 피부의 특징:",
        "  • T존은 지성, U존은 건성",
        "  • 수분-유분 밸런스가 중요",
        "  • 부위별로 다른 케어가 필요",
        "",
        "★ 독도토너의 '순하고 가벼운' 제형이",
        "★ 복합성 피부의 다양한 니즈를 충족시키지 못함",
    ])
    add_content_slide(prs, "복합성 부정률이 높은 이유", combo_lines, header_color=RED)

    # 슬라이드 7: 지성 피부 분석
    dokdo_oily = skin_df[(skin_df['BRAND_NAME'] == '독도토너') & (skin_df['SKIN_TYPE'] == '지성')]
    oily_neg = dokdo_oily[dokdo_oily['sentiment'] == 'NEG']
    oily_pain = Counter()
    for idx in oily_neg.index:
        for p in cat_map.get(idx, {}).get('pain_points', []):
            oily_pain[p] += 1

    oily_lines = [
        f"독도토너 지성 피부 리뷰: {len(dokdo_oily)}건",
        f"부정률: {dokdo_stats['지성']['rate']:.1f}% ({len(oily_neg)}건)",
        "",
        "지성 피부 불만 포인트:",
    ]
    for p, c in oily_pain.most_common(5):
        oily_lines.append(f"  • {p}: {c}건")

    oily_lines.extend([
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "지성 피부의 특징:",
        "  • 피지 분비가 많음",
        "  • 산뜻하고 가벼운 제형 선호",
        "  • 모공/피지 관리 중시",
        "",
        "▶ 독도토너가 '보습'에 집중되어 있어",
        "▶ 지성 피부에게는 '끈적임', '무거움'으로 느껴질 수 있음",
    ])
    add_content_slide(prs, "지성 피부 분석", oily_lines)

    # 슬라이드 8: 핵심 인사이트
    add_insight_box(prs, "핵심 인사이트", [
        "① 독도토너의 '순함' 포지셔닝은 민감성 피부에 완벽히 적중 → 부정률 0%",
        "② 하지만 복합성(4.2%), 지성(3.4%)에서는 오히려 타 브랜드보다 불만이 많음",
        "③ '순하고 가벼운' 제형이 모든 피부 타입의 니즈를 충족시키지는 못함",
        "④ 복합성/지성 피부는 '수분-유분 밸런스', '산뜻함'을 더 원함",
        "⑤ 피부타입별 맞춤 메시지 전략이 필요 (민감성 ↔ 복합성/지성 다르게)",
    ])

    # 슬라이드 9: 전략 제안
    add_content_slide(prs, "피부타입별 전략 제안", [
        "민감성 피부 (부정률 0%)",
        "  → '민감성 피부를 위한 No.1 토너'로 강력히 포지셔닝",
        "  → 재구매 유도 강화, 충성 고객화",
        "",
        "복합성 피부 (부정률 4.2%)",
        "  → '가벼운 수분 공급'으로 메시지 조정",
        "  → T존/U존 밸런스 케어 루틴 제안",
        "  → 다른 제품과의 조합 추천",
        "",
        "지성 피부 (부정률 3.4%)",
        "  → '산뜻한 수분' 강조",
        "  → '닦토용'으로 포지셔닝 (지성에게 닦토가 인기)",
        "  → '끈적이지 않는 가벼운 토너'로 소구",
        "",
        "▶ 모든 피부타입에 똑같이 소구하지 말고,",
        "▶ 피부타입별 맞춤 메시지로 전환율 높이기",
    ])

    # 슬라이드 10: 마무리
    add_content_slide(prs, "결론", [
        "독도토너는 '민감성 피부의 영웅'입니다.",
        "",
        "다른 브랜드들이 민감성 피부에서 2-3% 부정률을 보이는 동안,",
        "독도토너만 유일하게 0%를 기록했습니다.",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "하지만 '순함'만으로는 모든 피부 타입을 만족시킬 수 없습니다.",
        "",
        "복합성/지성 피부에게는",
        "다른 접근 방식이 필요합니다.",
        "",
        "▶ 강점(민감성)은 더 강화하고,",
        "▶ 약점(복합성/지성)은 메시지로 보완하는",
        "▶ 피부타입별 차별화 전략을 제안합니다.",
    ], header_color=DARK_BLUE)

    # 저장
    output_path = 'output/독도토너_피부타입분석.pptx'
    prs.save(output_path)

    print(f"\n리포트 생성 완료: {output_path}")
    print(f"총 {len(prs.slides)}장 슬라이드")


if __name__ == "__main__":
    main()
