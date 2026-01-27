# -*- coding: utf-8 -*-
"""
기존 PPT에 리뷰 분석 결과 추가
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import copy

# 기존 PPT 열기
prs = Presentation('docs/프로젝트2-이커머스1팀.pptx')

def add_content_slide(title, content_lines, insert_idx=None):
    """내용 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # blank

    if insert_idx is not None:
        # 특정 위치에 삽입
        slide = prs.slides.add_slide(slide_layout)
        # 슬라이드 순서 조정
        slide_id = prs.slides._sldIdLst[-1]
        prs.slides._sldIdLst.remove(slide_id)
        prs.slides._sldIdLst.insert(insert_idx, slide_id)
    else:
        slide = prs.slides.add_slide(slide_layout)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True

    # 내용
    content_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(9.2), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.space_after = Pt(6)

    return slide

def add_table_slide(title, headers, rows, insert_idx=None):
    """테이블 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]

    if insert_idx is not None:
        slide = prs.slides.add_slide(slide_layout)
        slide_id = prs.slides._sldIdLst[-1]
        prs.slides._sldIdLst.remove(slide_id)
        prs.slides._sldIdLst.insert(insert_idx, slide_id)
    else:
        slide = prs.slides.add_slide(slide_layout)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True

    # 테이블
    cols = len(headers)
    row_height = 0.35
    table = slide.shapes.add_table(
        len(rows) + 1, cols,
        Inches(0.3), Inches(0.8),
        Inches(9.4), Inches(row_height * (len(rows) + 1))
    ).table

    # 헤더
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)

    # 데이터
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(10)

    return slide

def add_image_slide(title, image_path, insert_idx=None):
    """이미지 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]

    if insert_idx is not None:
        slide = prs.slides.add_slide(slide_layout)
        slide_id = prs.slides._sldIdLst[-1]
        prs.slides._sldIdLst.remove(slide_id)
        prs.slides._sldIdLst.insert(insert_idx, slide_id)
    else:
        slide = prs.slides.add_slide(slide_layout)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True

    # 이미지
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.2), Inches(0.6), width=Inches(9.6))

    return slide

def add_two_images_slide(title, img1, img2, insert_idx=None):
    """이미지 2개 슬라이드"""
    slide_layout = prs.slide_layouts[6]

    if insert_idx is not None:
        slide = prs.slides.add_slide(slide_layout)
        slide_id = prs.slides._sldIdLst[-1]
        prs.slides._sldIdLst.remove(slide_id)
        prs.slides._sldIdLst.insert(insert_idx, slide_id)
    else:
        slide = prs.slides.add_slide(slide_layout)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(9.2), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True

    # 이미지 2개
    if os.path.exists(img1):
        slide.shapes.add_picture(img1, Inches(0.1), Inches(0.55), width=Inches(4.9))
    if os.path.exists(img2):
        slide.shapes.add_picture(img2, Inches(5.0), Inches(0.55), width=Inches(4.9))

    return slide

print("기존 PPT에 리뷰 분석 결과 추가 중...")
print(f"현재 슬라이드 수: {len(prs.slides)}")

# ===== 슬라이드 추가 =====

# 슬라이드 5 뒤에 (현상 진단 섹션에) 경쟁사 비교 데이터 추가
# 기존 슬라이드 7-8 사이에 삽입 (인덱스 7)

add_content_slide(
    "올리브영 토너 리뷰 27,745건 분석 (v2.0)",
    [
        "분석 대상: 7개 브랜드 토너 카테고리 전체 리뷰",
        "",
        "• 토리든: 8,065건 (29.1%)",
        "• 브링그린: 6,312건 (22.7%)",
        "• 독도토너: 3,750건 (13.5%)",
        "• 아누아: 3,544건 (12.8%)",
        "• 에스네이처: 2,232건 (8.0%)",
        "• 토니모리: 2,195건 (7.9%)",
        "• 아비브: 1,647건 (5.9%)",
        "",
        "분석 방법: 키워드 기반 감성 분석 + 역접 패턴/피부질병 탐지 로직 적용"
    ],
    insert_idx=7
)

add_table_slide(
    "브랜드별 실제 부정 비율 (별점 아닌 내용 기반)",
    ["브랜드", "리뷰수", "평점", "NEG비율", "피부질병 언급", "특징"],
    [
        ["독도토너", "3,750", "4.9", "11.1%", "3.3%", "중간 수준"],
        ["토리든", "8,065", "4.8", "10.3%", "3.7%", "보습 강자"],
        ["토니모리", "2,195", "4.9", "6.9%", "2.2%", "가장 낮음"],
        ["아누아", "3,544", "4.8", "21.3%", "13.0%", "진정 but 자극"],
        ["브링그린", "6,312", "4.8", "20.7%", "12.6%", "어성초 트러블"],
    ],
    insert_idx=8
)

# 포지셔닝 맵 추가 (슬라이드 9)
add_two_images_slide(
    "브랜드 포지셔닝: 효능(진정vs보습) & 사용감(물같음vs쫀쫀)",
    "output/figures/3B_positioning_benefit.png",
    "output/figures/3B_positioning_texture.png",
    insert_idx=9
)

# 원인 분석 섹션 (슬라이드 12 뒤) - Pain Point 추가
add_image_slide(
    "저평점 리뷰 Pain Point 분석",
    "output/figures/pain_point_heatmap.png",
    insert_idx=13
)

# 재구매/이탈 분석 (Retention 섹션에)
add_two_images_slide(
    "재구매 vs 이탈 신호 분석",
    "output/figures/loyalty_comparison.png",
    "output/figures/net_loyalty.png",
    insert_idx=17
)

# 독도토너 인사이트 정리
add_content_slide(
    "독도토너 리뷰 분석 인사이트",
    [
        "📊 현재 위치",
        "• NEG 비율 11.1%로 중간 수준 (토니모리 6.9% < 독도 < 아누아 21.3%)",
        "• 피부질병 언급률 3.3%로 양호 (진정 브랜드 대비 낮음)",
        "• 평균 평점 4.9점으로 최상위권",
        "",
        "💡 기회 요인",
        "• '무난함'이 약점이자 강점 - 자극 리스크 낮음",
        "• 진정 특화 브랜드(아누아, 브링그린)의 높은 부작용 대비 안전성 어필 가능",
        "",
        "⚠️ 위협 요인",
        "• 토리든이 보습 포지션에서 압도적 리뷰 수(8,065건) 확보",
        "• '무난/애매' 태그 비율 높음 → 차별화 포인트 부재",
        "",
        "🎯 제언",
        "• '순함'에서 '피부장벽 강화' 메시지로 전환",
        "• 경쟁사 대비 낮은 부작용률 데이터 마케팅 활용"
    ],
    insert_idx=18
)

# 전환 매트릭스
add_image_slide(
    "고객 전환 패턴 (브랜드 간 이동)",
    "output/figures/3E_switch_matrix.png",
    insert_idx=19
)

# 종합 대시보드
add_image_slide(
    "리뷰 분석 종합 대시보드",
    "output/figures/dashboard_summary.png",
    insert_idx=20
)

print(f"추가 후 슬라이드 수: {len(prs.slides)}")

# 저장
output_path = "docs/프로젝트2-이커머스1팀_분석추가.pptx"
prs.save(output_path)
print(f"저장 완료: {output_path}")
