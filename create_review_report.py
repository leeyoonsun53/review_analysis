# -*- coding: utf-8 -*-
"""
토너 리뷰 분석 결과 리포트 생성
리뷰 데이터 분석 결과를 알기 쉽게 정리한 슬라이드
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import pandas as pd
import json
import sys
from datetime import datetime
from collections import Counter

sys.stdout.reconfigure(encoding='utf-8')

# 색상 정의
NAVY = RGBColor(0x1a, 0x36, 0x5d)
LIGHT_BLUE = RGBColor(0x4a, 0x90, 0xd9)
GREEN = RGBColor(0x27, 0xae, 0x60)
RED = RGBColor(0xe7, 0x4c, 0x3c)
ORANGE = RGBColor(0xf3, 0x9c, 0x12)
GRAY = RGBColor(0x7f, 0x8c, 0x8d)
WHITE = RGBColor(0xff, 0xff, 0xff)
DARK_GRAY = RGBColor(0x2c, 0x3e, 0x50)
LIGHT_GRAY = RGBColor(0xf8, 0xf9, 0xfa)


def add_title_slide(prs, title, subtitle):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 상단 장식
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()

    # 배경 박스
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), Inches(10), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_section_slide(prs, section_num, title):
    """섹션 구분 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(10), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()

    # 섹션 번호
    if section_num:
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(9), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_num
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, content_lines, highlight_indices=None):
    """내용 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 내용
    content_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.1), Inches(9.2), Inches(6))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

        if highlight_indices and i in highlight_indices:
            p.font.bold = True
            p.font.color.rgb = LIGHT_BLUE


def add_table_slide(prs, title, headers, rows, subtitle=""):
    """테이블 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 서브타이틀
    start_y = 1.1
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(9.2), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        start_y = 1.4

    # 테이블
    num_rows = len(rows) + 1
    num_cols = len(headers)
    row_height = min(0.4, 5.5 / num_rows)

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.3), Inches(start_y),
        Inches(9.4), Inches(row_height * num_rows)
    ).table

    # 헤더
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # 데이터
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT


def add_bar_chart_slide(prs, title, categories, series_data, subtitle=""):
    """막대 차트 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 서브타이틀
    start_y = 1.0
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(9.2), Inches(0.3))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        start_y = 1.3

    # 차트 데이터
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series_name, values in series_data.items():
        chart_data.add_series(series_name, values)

    # 차트 추가
    x, y, cx, cy = Inches(0.5), Inches(start_y), Inches(9), Inches(5.5)
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)


def add_insight_slide(prs, title, insights):
    """인사이트 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더 (오렌지)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 인사이트 박스들
    y_pos = 1.2
    for insight in insights:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3), Inches(y_pos),
            Inches(9.4), Inches(0.85)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = LIGHT_BLUE

        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.1), Inches(9), Inches(0.7))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = insight
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY

        y_pos += 0.95


def add_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    """2단 컬럼 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 왼쪽 컬럼 - 초록 헤더
    left_header = slide.shapes.add_textbox(Inches(0.4), Inches(1.1), Inches(4.3), Inches(0.4))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN

    left_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.5), Inches(4.3), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

    # 오른쪽 컬럼 - 빨강 헤더
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RED

    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)


def add_quote_slide(prs, quote):
    """강조 인용 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경 박스
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(9), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = LIGHT_BLUE
    shape.line.width = Pt(3)

    # 인용문
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(8), Inches(3))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER


def main():
    print("=" * 60)
    print("토너 리뷰 분석 결과 리포트 생성")
    print("=" * 60)

    # ========================================
    # 데이터 로드
    # ========================================
    print("\n[데이터 로드 중...]")
    df = pd.read_csv('data/merged_reviews_processed.csv', encoding='utf-8-sig')
    cat = json.load(open('output/gpt_analysis_categorized.json', encoding='utf-8'))
    points = json.load(open('output/points_categorized.json', encoding='utf-8'))

    # 기본 통계 계산
    total_reviews = len(df)
    platforms = df['PLATFORM'].unique()
    brands = ['토리든', '브링그린', '독도토너', '에스네이처', '아누아', '토니모리', '아비브']

    pos_count = sum(1 for r in cat if r.get('sentiment') == 'POS')
    neg_count = sum(1 for r in cat if r.get('sentiment') == 'NEG')
    neu_count = total_reviews - pos_count - neg_count
    pos_rate = pos_count / total_reviews * 100
    neg_rate = neg_count / total_reviews * 100
    neu_rate = neu_count / total_reviews * 100

    print(f"   총 리뷰: {total_reviews:,}건")
    print(f"   플랫폼: {', '.join(platforms)}")
    print(f"   긍정: {pos_count:,}건 ({pos_rate:.1f}%)")
    print(f"   부정: {neg_count:,}건 ({neg_rate:.1f}%)")

    # 브랜드별 통계
    brand_stats = {}
    for brand in brands:
        brand_cat = [r for r in cat if r.get('brand') == brand]
        total = len(brand_cat)
        pos = sum(1 for r in brand_cat if r.get('sentiment') == 'POS')
        neg = sum(1 for r in brand_cat if r.get('sentiment') == 'NEG')
        brand_stats[brand] = {
            'total': total,
            'pos': pos,
            'neg': neg,
            'pos_rate': pos / total * 100 if total > 0 else 0,
            'neg_rate': neg / total * 100 if total > 0 else 0
        }

    # Pain/Positive/Usage 통계
    pain_counts = Counter()
    pos_counts = Counter()
    usage_counts = Counter()
    benefit_counts = Counter()
    texture_counts = Counter()

    for r in cat:
        for pain in r.get('pain_points', []):
            pain_counts[pain] += 1
        for pos in r.get('positive_points', []):
            pos_counts[pos] += 1
        for u in r.get('usage_tags', []):
            usage_counts[u] += 1
        for b in r.get('benefit_tags', []):
            benefit_counts[b] += 1
        for t in r.get('texture_tags', []):
            texture_counts[t] += 1

    # 브랜드별 usage
    brand_usage = {}
    for brand in brands:
        brand_cat = [r for r in cat if r.get('brand') == brand]
        brand_total = len(brand_cat)
        usage_stats = Counter()
        for r in brand_cat:
            for u in r.get('usage_tags', []):
                usage_stats[u] += 1
        brand_usage[brand] = {
            '닦토': usage_stats.get('닦토', 0) / brand_total * 100 if brand_total > 0 else 0,
            '레이어링': usage_stats.get('레이어링', 0) / brand_total * 100 if brand_total > 0 else 0,
            '스킨팩/토너팩': usage_stats.get('스킨팩/토너팩', 0) / brand_total * 100 if brand_total > 0 else 0,
        }

    # 브랜드별 benefit/texture
    brand_benefit = {}
    brand_texture = {}
    for brand in brands:
        brand_cat = [r for r in cat if r.get('brand') == brand]
        brand_total = len(brand_cat)
        b_stats = Counter()
        t_stats = Counter()
        for r in brand_cat:
            for b in r.get('benefit_tags', []):
                b_stats[b] += 1
            for t in r.get('texture_tags', []):
                t_stats[t] += 1
        brand_benefit[brand] = {k: v/brand_total*100 if brand_total > 0 else 0 for k, v in b_stats.items()}
        brand_texture[brand] = {k: v/brand_total*100 if brand_total > 0 else 0 for k, v in t_stats.items()}

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========================================
    # 슬라이드 1: 표지
    # ========================================
    print("\n[슬라이드 생성 중...]")
    add_title_slide(prs,
        "토너 리뷰 분석 리포트",
        f"7개 브랜드 {total_reviews:,}건 리뷰 분석 | GPT-4o-mini 기반 | {datetime.now().strftime('%Y.%m')}")

    # ========================================
    # 슬라이드 2: 분석 개요
    # ========================================
    platform_counts = df['PLATFORM'].value_counts()
    overview_lines = [
        "분석 개요",
        "",
        f"  총 리뷰 수: {total_reviews:,}건",
        f"  분석 기간: 2025.02 ~ 2026.01",
        f"  분석 방법: GPT-4o-mini 기반 감성분석 + 키워드 추출",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "플랫폼별 리뷰 수",
    ]
    for platform, count in platform_counts.items():
        overview_lines.append(f"  • {platform}: {count:,}건 ({count/total_reviews*100:.1f}%)")

    overview_lines.extend(["", "분석 브랜드 (7개)"])
    for brand in brands:
        stats = brand_stats[brand]
        overview_lines.append(f"  • {brand}: {stats['total']:,}건")

    add_content_slide(prs, "분석 개요", overview_lines, highlight_indices=[0, 2, 3, 4])

    # ========================================
    # 슬라이드 3: 감성 분석 결과
    # ========================================
    add_section_slide(prs, "PART 1", "전체 감성 분석")

    sentiment_lines = [
        "전체 감성 분포",
        "",
        f"  긍정(POS): {pos_count:,}건 ({pos_rate:.1f}%)",
        f"  중립(NEU): {neu_count:,}건 ({neu_rate:.1f}%)",
        f"  부정(NEG): {neg_count:,}건 ({neg_rate:.1f}%)",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "해석",
        "",
        f"  • 전체 긍정률 {pos_rate:.1f}%로 매우 높은 만족도",
        f"  • 부정 리뷰 {neg_rate:.1f}%는 개선 포인트 파악에 집중",
        f"  • 중립 리뷰 {neu_rate:.1f}%는 '무난하다'는 평가 다수",
    ]
    add_content_slide(prs, "전체 감성 분석 결과", sentiment_lines, highlight_indices=[2, 3, 4])

    # ========================================
    # 슬라이드 4: 브랜드별 감성 비교
    # ========================================
    brand_rows = []
    for brand in sorted(brands, key=lambda x: brand_stats[x]['pos_rate'], reverse=True):
        stats = brand_stats[brand]
        brand_rows.append([
            brand,
            f"{stats['total']:,}",
            f"{stats['pos_rate']:.1f}%",
            f"{stats['neg_rate']:.1f}%"
        ])

    add_table_slide(prs, "브랜드별 감성 비교",
        ["브랜드", "리뷰 수", "긍정률", "부정률"],
        brand_rows,
        subtitle="긍정률 기준 정렬")

    # 브랜드별 긍정률 차트
    sorted_brands = sorted(brands, key=lambda x: brand_stats[x]['pos_rate'], reverse=True)
    add_bar_chart_slide(prs, "브랜드별 긍정률 비교",
        categories=sorted_brands,
        series_data={
            '긍정률(%)': [brand_stats[b]['pos_rate'] for b in sorted_brands],
            '부정률(%)': [brand_stats[b]['neg_rate'] for b in sorted_brands]
        },
        subtitle="에스네이처 > 토니모리 > 아비브 > 독도토너 순")

    # ========================================
    # 슬라이드 5: Positive Points 분석
    # ========================================
    add_section_slide(prs, "PART 2", "고객이 좋아하는 포인트")

    top_positives = pos_counts.most_common(12)
    pos_lines = ["TOP 12 Positive Points", ""]
    for i, (pos, cnt) in enumerate(top_positives, 1):
        pos_lines.append(f"  {i:2d}. {pos}: {cnt:,}건")

    pos_lines.extend([
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "핵심 인사이트",
        "  • '순함/저자극'이 압도적 1위 - 토너 선택 핵심 기준",
        "  • '촉촉함/보습'이 2위 - 기본 효능 중시",
        "  • '가성비/대용량' 다수 - 실용적 가치 중요",
    ])
    add_content_slide(prs, "고객이 좋아하는 포인트 (Positive Points)", pos_lines, highlight_indices=[2, 3, 4])

    # Positive 차트
    top_pos_names = [p[0][:8] for p in top_positives[:8]]  # 이름 축약
    top_pos_counts = [p[1] for p in top_positives[:8]]
    add_bar_chart_slide(prs, "Positive Points TOP 8",
        categories=top_pos_names,
        series_data={'건수': top_pos_counts})

    # ========================================
    # 슬라이드 6: Pain Points 분석
    # ========================================
    add_section_slide(prs, "PART 3", "고객 불만 포인트")

    top_pains = pain_counts.most_common(12)
    pain_lines = ["TOP 12 Pain Points", ""]
    for i, (pain, cnt) in enumerate(top_pains, 1):
        pain_lines.append(f"  {i:2d}. {pain}: {cnt:,}건")

    pain_lines.extend([
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "핵심 인사이트",
        "  • '건조/보습부족' 1위 - 전 브랜드 공통 과제",
        "  • '트러블/뾰루지' 다수 - 피부 반응 우려",
        "  • '효과 없음' - 기대 대비 실망감",
    ])
    add_content_slide(prs, "고객 불만 포인트 (Pain Points)", pain_lines, highlight_indices=[2, 3, 4])

    # Pain 카테고리 분포
    pain_cat = points.get('pain_points', {}).get('category_counts', {})
    pain_cat_sorted = sorted([(k, v) for k, v in pain_cat.items() if k != '기타'], key=lambda x: x[1], reverse=True)[:8]
    add_bar_chart_slide(prs, "Pain Point 카테고리 분포",
        categories=[p[0] for p in pain_cat_sorted],
        series_data={'건수': [p[1] for p in pain_cat_sorted]},
        subtitle="자극/트러블, 건조/보습부족이 주요 불만 요인")

    # ========================================
    # 슬라이드 7: 브랜드별 강점/약점
    # ========================================
    add_section_slide(prs, "PART 4", "브랜드별 분석")

    # 브랜드별 효능 포지셔닝
    benefit_rows = []
    for brand in brands:
        b = brand_benefit.get(brand, {})
        benefit_rows.append([
            brand,
            f"{b.get('보습', 0):.1f}%",
            f"{b.get('진정', 0):.1f}%",
            f"{b.get('장벽', 0):.1f}%",
            f"{b.get('결', 0):.1f}%"
        ])

    add_table_slide(prs, "브랜드별 효능 태그 비율",
        ["브랜드", "보습", "진정", "장벽", "결 개선"],
        benefit_rows,
        subtitle="리뷰에서 추출한 효능 키워드 언급 비율")

    # 브랜드별 사용감 포지셔닝
    texture_rows = []
    for brand in brands:
        t = brand_texture.get(brand, {})
        texture_rows.append([
            brand,
            f"{t.get('물같음', 0):.1f}%",
            f"{t.get('쫀쫀', 0):.1f}%",
            f"{t.get('끈적', 0):.1f}%",
            f"{t.get('흡수', 0):.1f}%"
        ])

    add_table_slide(prs, "브랜드별 사용감 태그 비율",
        ["브랜드", "물같음", "쫀쫀", "끈적임", "흡수력"],
        texture_rows,
        subtitle="리뷰에서 추출한 사용감 키워드 언급 비율")

    # ========================================
    # 슬라이드 8: 사용법 분석
    # ========================================
    add_section_slide(prs, "PART 5", "사용법(Usage) 분석")

    top_usages = usage_counts.most_common(6)
    usage_lines = ["고객이 언급한 사용법 TOP 6", ""]
    for i, (usage, cnt) in enumerate(top_usages, 1):
        pct = cnt / total_reviews * 100
        usage_lines.append(f"  {i}. {usage}: {cnt:,}건 ({pct:.1f}%)")

    usage_lines.extend([
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "핵심 인사이트",
        "  • '닦토(닦아내는 토너)' 사용이 가장 많음",
        "  • 레이어링(7스킨법 등) 사용도 다수",
        "  • 스킨팩/토너팩으로 활용하는 고객 존재",
        "",
        "마케팅 활용",
        "  → 사용법별 맞춤 메시지 가능",
        "  → 대용량/가성비 + 닦토 조합 강조",
    ])
    add_content_slide(prs, "사용법(Usage) 분석", usage_lines, highlight_indices=[2, 3])

    # 브랜드별 사용법 비교
    usage_rows = []
    for brand in brands:
        u = brand_usage.get(brand, {})
        usage_rows.append([
            brand,
            f"{u.get('닦토', 0):.1f}%",
            f"{u.get('레이어링', 0):.1f}%",
            f"{u.get('스킨팩/토너팩', 0):.1f}%"
        ])

    add_table_slide(prs, "브랜드별 사용법 비율",
        ["브랜드", "닦토", "레이어링", "스킨팩"],
        usage_rows,
        subtitle="독도토너/토니모리 = 닦토용 가성비템으로 포지셔닝")

    # ========================================
    # 슬라이드 9: 독도토너 심층 분석
    # ========================================
    add_section_slide(prs, "PART 6", "독도토너 심층 분석")

    # 독도토너 데이터 추출
    dokdo_cat = [r for r in cat if r.get('brand') == '독도토너']
    dokdo_total = len(dokdo_cat)
    dokdo_pos = sum(1 for r in dokdo_cat if r.get('sentiment') == 'POS')
    dokdo_neg = sum(1 for r in dokdo_cat if r.get('sentiment') == 'NEG')

    dokdo_pain = Counter()
    dokdo_positive = Counter()
    for r in dokdo_cat:
        for p in r.get('pain_points', []):
            dokdo_pain[p] += 1
        for p in r.get('positive_points', []):
            dokdo_positive[p] += 1

    # 독도토너 강점/약점
    top_dokdo_pos = dokdo_positive.most_common(5)
    top_dokdo_pain = dokdo_pain.most_common(5)

    add_two_column_slide(prs, "독도토너 강점 vs 약점",
        "강점 (Positive Points)",
        [f"{i}. {p[0]}: {p[1]}건" for i, p in enumerate(top_dokdo_pos, 1)] + [
            "",
            "→ '순함/저자극' 압도적 1위",
            "→ '촉촉함', '무난함' 다수",
            "→ 기본기가 탄탄한 토너"
        ],
        "약점 (Pain Points)",
        [f"{i}. {p[0]}: {p[1]}건" for i, p in enumerate(top_dokdo_pain, 1)] + [
            "",
            "→ '건조', '효과부족' 언급",
            "→ 저평점 중 46%가 '효과 없음'",
            "→ '순함'만으로는 부족"
        ])

    # 독도토너 상세 분석
    dokdo_lines = [
        f"분석 대상: {dokdo_total:,}건 리뷰",
        "",
        f"  긍정: {dokdo_pos:,}건 ({dokdo_pos/dokdo_total*100:.1f}%)",
        f"  부정: {dokdo_neg:,}건 ({dokdo_neg/dokdo_total*100:.1f}%)",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "독도토너 포지셔닝",
        "",
        "  현재 인식:",
        "  • '무난한 입문용 토너'",
        "  • '순하고 자극 없는 토너'",
        "  • '닦토용 가성비템'",
        "",
        "  문제점:",
        "  • '순함'만으로는 재구매 동기 부족",
        "  • '효과 없음' 인식 → 다른 브랜드로 이탈",
        "  • Nice to have → Must have 전환 필요",
    ]
    add_content_slide(prs, "독도토너 현황 분석", dokdo_lines, highlight_indices=[0, 2, 3, 15, 16, 17])

    # ========================================
    # 슬라이드 10: 핵심 인사이트
    # ========================================
    add_section_slide(prs, "SUMMARY", "핵심 인사이트")

    add_insight_slide(prs, "리뷰 분석 핵심 인사이트", [
        f"① 전체 긍정률 {pos_rate:.1f}% - 토너 시장 전반적 만족도 높음",
        "② '순함/저자극'이 핵심 구매 기준 - 5,000건 이상 언급",
        "③ '건조/보습부족'이 최대 불만 - 전 브랜드 공통 과제",
        "④ '닦토' 사용이 가장 많음 - 대용량/가성비 중요",
        f"⑤ 독도토너는 '순함'에서 강점, '효과부족'이 약점",
    ])

    # 브랜드별 핵심 특징
    brand_features = [
        "토리든: 보습 강자, 건조 불만도 1위 (양날의 검)",
        "브링그린: 진정 효과 좋으나 트러블/자극 불만 존재",
        "독도토너: 순함의 대명사, 효과부족 인식 개선 필요",
        "에스네이처: 긍정률 1위, 안정적 품질",
        "아누아: 프리미엄 진정, 가격 불만",
        "토니모리: 가성비 1위, 닦토용으로 인기",
        "아비브: 진정 효과, 프리미엄 이미지",
    ]
    add_content_slide(prs, "브랜드별 핵심 특징", [""] + brand_features)

    # ========================================
    # 슬라이드 11: 마무리
    # ========================================
    add_quote_slide(prs,
        f"총 {total_reviews:,}건의 리뷰에서\n"
        "고객의 진짜 목소리를 분석했습니다.\n\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n\n"
        "'순함'은 기본, '효과'로 차별화\n"
        "'가성비'는 유지, '가치'를 더하기\n\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n\n"
        "감사합니다.")

    # ========================================
    # 저장
    # ========================================
    output_path = 'output/토너_리뷰분석_결과보고서.pptx'
    prs.save(output_path)

    print("\n" + "=" * 60)
    print(f"리포트 생성 완료: {output_path}")
    print(f"총 {len(prs.slides)}장 슬라이드")
    print("=" * 60)


if __name__ == "__main__":
    main()
