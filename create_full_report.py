# -*- coding: utf-8 -*-
"""
독도토너 분석 전체 리포트 생성 (55장)
PDF + PPT 통합 버전
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import pandas as pd
import json
import sys
from datetime import datetime
from collections import Counter

sys.stdout.reconfigure(encoding='utf-8')

# 색상 정의
NAVY = RGBColor(0x1a, 0x36, 0x5d)
LIGHT_BLUE = RGBColor(0x4a, 0x90, 0xd9)
GREEN = RGBColor(0x27, 0xae, 0x60)
RED = RGBColor(0xe7, 0x4c, 0x3c)
ORANGE = RGBColor(0xf3, 0x9c, 0x12)
GRAY = RGBColor(0x7f, 0x8c, 0x8d)
WHITE = RGBColor(0xff, 0xff, 0xff)
DARK_GRAY = RGBColor(0x2c, 0x3e, 0x50)


def add_title_slide(prs, title, subtitle):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), Inches(10), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_section_slide(prs, section_num, title, subtitle=""):
    """섹션 구분 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(10), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()

    # 섹션 번호
    if section_num:
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(9), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_num
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.1) if section_num else Inches(2.8), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, content_lines, highlight_indices=None):
    """내용 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 내용
    content_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.1), Inches(9.2), Inches(6))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

        if highlight_indices and i in highlight_indices:
            p.font.bold = True
            p.font.color.rgb = LIGHT_BLUE


def add_table_slide(prs, title, headers, rows, subtitle=""):
    """테이블 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 서브타이틀
    start_y = 1.1
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(9.2), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        start_y = 1.4

    # 테이블
    num_rows = len(rows) + 1
    num_cols = len(headers)
    row_height = min(0.4, 5.5 / num_rows)

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.3), Inches(start_y),
        Inches(9.4), Inches(row_height * num_rows)
    ).table

    # 헤더
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # 데이터
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT


def add_insight_slide(prs, title, insights):
    """인사이트 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더 (오렌지)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 인사이트 박스들
    y_pos = 1.2
    for insight in insights:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3), Inches(y_pos),
            Inches(9.4), Inches(0.85)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xf8, 0xf9, 0xfa)
        box.line.color.rgb = LIGHT_BLUE

        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.1), Inches(9), Inches(0.7))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = insight
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY

        y_pos += 0.95


def add_quote_slide(prs, quote, source=""):
    """강조 인용 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경 박스
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(9), Inches(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xf8, 0xf9, 0xfa)
    shape.line.color.rgb = LIGHT_BLUE
    shape.line.width = Pt(3)

    # 인용문
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    if source:
        src_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(0.5))
        tf = src_box.text_frame
        p = tf.paragraphs[0]
        p.text = source
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER


def add_bar_chart_slide(prs, title, categories, series_data, subtitle=""):
    """막대 차트 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 서브타이틀
    start_y = 1.0
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(9.2), Inches(0.3))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        start_y = 1.3

    # 차트 데이터
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series_name, values in series_data.items():
        chart_data.add_series(series_name, values)

    # 차트 추가
    x, y, cx, cy = Inches(0.5), Inches(start_y), Inches(9), Inches(5.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    return slide


def add_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    """2단 컬럼 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 왼쪽 컬럼
    left_header = slide.shapes.add_textbox(Inches(0.4), Inches(1.1), Inches(4.3), Inches(0.4))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

    left_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.5), Inches(4.3), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

    # 오른쪽 컬럼
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)


def main():
    print("=" * 60)
    print("독도토너 분석 전체 리포트 생성 (55장)")
    print("=" * 60)

    # ========================================
    # 데이터 로드
    # ========================================
    print("\n[데이터 로드 중...]")
    df = pd.read_csv('data/merged_reviews_processed.csv', encoding='utf-8-sig')
    cat = json.load(open('output/gpt_analysis_categorized.json', encoding='utf-8'))

    # 기본 통계 계산
    total_reviews = len(df)
    pos_count = sum(1 for r in cat if r.get('sentiment') == 'POS')
    neg_count = sum(1 for r in cat if r.get('sentiment') == 'NEG')
    neu_count = total_reviews - pos_count - neg_count
    pos_rate = pos_count / total_reviews * 100
    neg_rate = neg_count / total_reviews * 100

    print(f"   총 리뷰: {total_reviews:,}건")
    print(f"   긍정: {pos_count:,}건 ({pos_rate:.1f}%)")
    print(f"   부정: {neg_count:,}건 ({neg_rate:.1f}%)")

    # 브랜드별 통계
    brands = ['토리든', '브링그린', '독도토너', '에스네이처', '아누아', '토니모리', '아비브']
    brand_stats = {}
    for brand in brands:
        brand_cat = [r for r in cat if r.get('brand') == brand]
        total = len(brand_cat)
        pos = sum(1 for r in brand_cat if r.get('sentiment') == 'POS')
        neg = sum(1 for r in brand_cat if r.get('sentiment') == 'NEG')
        brand_stats[brand] = {
            'total': total,
            'pos': pos,
            'neg': neg,
            'pos_rate': pos / total * 100 if total > 0 else 0,
            'neg_rate': neg / total * 100 if total > 0 else 0
        }

    # Pain Points 통계 계산
    pain_counts = Counter()
    for r in cat:
        for pain in r.get('pain_points', []):
            pain_counts[pain] += 1

    # Positive Points 통계 계산
    pos_counts = Counter()
    for r in cat:
        for pos in r.get('positive_points', []):
            pos_counts[pos] += 1

    # Usage 통계 계산
    usage_counts = Counter()
    for r in cat:
        for u in r.get('usage_tags', []):
            usage_counts[u] += 1

    # 브랜드별 사용법 비율 계산
    brand_usage = {}
    for brand in brands:
        brand_cat = [r for r in cat if r.get('brand') == brand]
        brand_total = len(brand_cat)
        usage_stats = Counter()
        for r in brand_cat:
            for u in r.get('usage_tags', []):
                usage_stats[u] += 1
        brand_usage[brand] = {
            '닦토': usage_stats.get('닦토', 0) / brand_total * 100 if brand_total > 0 else 0,
            '레이어링': usage_stats.get('레이어링', 0) / brand_total * 100 if brand_total > 0 else 0,
            '스킨팩/토너팩': usage_stats.get('스킨팩/토너팩', 0) / brand_total * 100 if brand_total > 0 else 0,
            '바디 사용': usage_stats.get('바디 사용', 0) / brand_total * 100 if brand_total > 0 else 0,
        }

    print(f"   Pain Points: {len(pain_counts)}종")
    print(f"   Positive Points: {len(pos_counts)}종")
    print(f"   Usage Tags: {len(usage_counts)}종")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========================================
    # Part 0. 표지 & 목차 (2장)
    # ========================================
    print("\n[Part 0] 표지 & 목차...")

    # 슬라이드 1: 표지
    add_title_slide(prs,
        "잊혀지는 국민 토너, 독도토너의 재도약 전략",
        f"라운드랩 자사몰 데이터 기반 Pre-Retention 분석 | 이커머스 1팀 | {datetime.now().strftime('%Y.%m')}")

    # 슬라이드 2: 목차/분석 프레임
    add_content_slide(prs, "분석 프레임워크", [
        "1. 현상 진단 - 잊혀지는 국민 토너의 현주소",
        "2. 원인 분석 - 데이터가 말하는 고객 이탈의 진실",
        "3. 전략 방향 - '무자극'을 넘어 '필수 베이스'로",
        "4. 실행 전술 - 데이터 기반의 맞춤형 솔루션",
        "5. 기대 효과 - KPI 및 성과 관리",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "데이터 소스:",
        "• 정량: 캐글 쇼핑몰 데이터 (자사몰 행동 로그)",
        "• 정성: 올리브영 27,745건 + 무신사 6,163건 = 33,908건 리뷰",
        "• 분석: GPT-4o-mini 기반 감성분석 (Pain/Positive/Usage 추출)",
    ], highlight_indices=[0, 1, 2, 3, 4])

    # ========================================
    # Part 1. 현상 진단 (10장)
    # ========================================
    print("[Part 1] 현상 진단...")

    # 슬라이드 3: 섹션 타이틀
    add_section_slide(prs, "1. 현상 진단", "잊혀지는 국민 토너의 현주소")

    # 슬라이드 4: 프로젝트 배경
    add_content_slide(prs, "프로젝트 배경", [
        "한때는 브랜드를 견인했던 '독도토너'",
        "→ 현재는 잊혀지는 추세",
        "→ 구매 활성화를 위한 전략 도출 필요",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "시장 변화:",
        "• 순한 라인 → 고기능 시대로 변화",
        "• 피부 장벽 니즈 세분화 → 리포지셔닝 필요",
        "• 재구매 가능성 있는 고객을 만드는 것이 핵심",
        "",
        "네이버 자사몰 전략:",
        "• UI/UX 개선을 통한 전환율 향상",
        "• CRM을 통한 재구매 유도",
    ], highlight_indices=[0, 1, 2])

    # 슬라이드 5: 프로젝트 목표
    add_content_slide(prs, "프로젝트 목표 (Goal)", [
        "독도토너 첫 구매 고객 중",
        "'재구매 가능 상태(Pre-Retention)'에 도달한 고객을 정의하고,",
        "해당 상태로의 전이를 가로막는 행동 병목을 규명한다.",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "목표를 한 줄로 요약하면:",
        "",
        "재구매 '결과'가 아닌,",
        "재구매로 이어질 '상태'를 데이터로 설명하는 것",
    ], highlight_indices=[0, 1, 2, 8, 9])

    # 슬라이드 6: RFM 재해석
    add_table_slide(prs, "핵심 분석 관점 - RFM 프레임 재해석",
        ["구분", "기존 의미", "본 프로젝트 해석"],
        [
            ["Recency", "최근 구매 시점", "구매 이후 관계가 유지되고 있는가?"],
            ["Frequency", "구매 빈도", "구매 외 반복 행동이 발생하고 있는가?"],
            ["Monetary", "구매 금액", "제품/브랜드에 대한 신뢰 확장 의지가 있는가?"],
        ],
        subtitle="→ RFM을 '고객 분류 도구'가 아닌 '상태 정의 도구'로 활용")

    # 슬라이드 7: 검색량 추이
    add_content_slide(prs, "라운드랩 vs 독도토너 - 네이버 검색량 추이", [
        "[네이버 검색 트렌드 차트 - 2017~2025]",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "▶ 브랜드 론칭 초기 (2017~2020)",
        "   - 독도 토너 검색량 > 라운드랩 검색량",
        "   - 독도 토너가 브랜드를 견인",
        "",
        "▶ 전환점 (2022년 이후)",
        "   - 라운드랩 ≈ 독도 토너 (검색량 동등)",
        "   - 독도 토너의 브랜드 견인력 상실",
        "",
        "▶ 결론",
        "   독도 토너의 브랜드 견인은 끝났다.",
        "   새로운 성장 동력이 필요한 시점",
    ], highlight_indices=[13, 14])

    # 슬라이드 8: 리뷰 분석 개요 (동적 데이터)
    brand_table_rows = []
    brand_features = {
        '토리든': '보습 강자', '브링그린': '진정 but 자극', '독도토너': '순함 중심',
        '에스네이처': '안정적', '아누아': '가격 불만', '토니모리': '가성비 리더', '아비브': '진정 효과'
    }
    for brand in brands:
        stats = brand_stats[brand]
        brand_table_rows.append([
            brand,
            f"{stats['total']:,}",
            f"{stats['total']/total_reviews*100:.1f}%",
            f"{stats['pos_rate']:.1f}%",
            f"{stats['neg_rate']:.1f}%",
            brand_features.get(brand, '')
        ])

    add_table_slide(prs, "올리브영 + 무신사 토너 리뷰 분석",
        ["브랜드", "리뷰수", "비중", "긍정률", "부정률", "특징"],
        brand_table_rows,
        subtitle=f"분석 대상: 7개 브랜드 총 {total_reviews:,}건 (2025.02~2026.01) | GPT-4o-mini 기반 감성분석")

    # 슬라이드 8-1: 브랜드별 긍정률 차트
    add_bar_chart_slide(prs, "브랜드별 긍정률 비교",
        categories=brands,
        series_data={
            '긍정률(%)': [brand_stats[b]['pos_rate'] for b in brands],
            '부정률(%)': [brand_stats[b]['neg_rate'] for b in brands]
        },
        subtitle="→ 에스네이처 1위, 독도토너 4위 (경쟁력 있는 포지션)")

    # 슬라이드 9: 브랜드 포지셔닝 맵
    add_two_column_slide(prs, "브랜드 포지셔닝 맵",
        "효능 축 (진정 vs 보습)",
        [
            "[보습 중심]",
            "• 토니모리 (보습 80.7%)",
            "• 에스네이처 (보습 73.4%)",
            "• 토리든 (보습 70.3%)",
            "",
            "[진정 중심]",
            "• 브링그린 (진정 59.7%)",
            "• 아비브 (진정 61.3%)",
            "• 아누아 (진정 57.3%)",
            "",
            "[균형형]",
            "• 독도토너 (보습 54.9% / 진정 52.5%)",
        ],
        "가치 축 (가성비 vs 프리미엄)",
        [
            "[가성비 리더]",
            "• 토니모리 (가성비 39.3%)",
            "• 브링그린 (가성비 21.8%)",
            "• 독도토너 (가성비 18.5%)",
            "",
            "[프리미엄 이미지]",
            "• 아누아 (인생템 60.3%)",
            "• 아비브 (인생템 59.3%)",
            "",
            "━━━━━━━━━━━━━━━━━━━",
            "독도토너 문제:",
            "명확한 정체성 부재",
            "→ '무난한 선택지'로 전락",
        ])

    # 슬라이드 10: GPT 감성분석 결과
    add_content_slide(prs, "GPT 감성분석 핵심 결과", [
        "▶ 전체 긍정률: 90.4% (부정률 2.0%)",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "항목별 1위 vs 독도토너:",
        "",
        "• 긍정률 1위: 에스네이처 (93.3%) → 독도토너 4위 (90.2%)",
        "• Positive Point 1위: 순함/저자극 (5,495건) → 독도토너 기여 2위",
        "• Pain Point 1위: 건조/보습부족 (463건) → 독도토너 기여 5위",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "▶ 독도토너 강점: 순함/저자극 1,062건 (브랜드 내 1위)",
        "▶ 독도토너 약점: '효과 없음' 인식 46% (저평점 리뷰 중)",
    ], highlight_indices=[0, 12, 13])

    # 슬라이드 11: Positive Points TOP 10 (동적)
    top_positives = pos_counts.most_common(10)
    pos_rows = []
    for i, (pos, cnt) in enumerate(top_positives, 1):
        pos_rows.append([str(i), pos, f"{cnt:,}", ""])

    add_table_slide(prs, "Positive Points TOP 10 (GPT 분석)",
        ["순위", "Positive Point", "건수", "주요 브랜드"],
        pos_rows,
        subtitle="→ 독도토너: '순함' + '무난함' = 기본기는 좋으나 WOW 요소 부족")

    # 슬라이드 11-1: Positive Points 차트
    top_pos_names = [p[0] for p in top_positives[:6]]
    top_pos_counts = [p[1] for p in top_positives[:6]]
    add_bar_chart_slide(prs, "Positive Points 분포 (상위 6개)",
        categories=top_pos_names,
        series_data={'건수': top_pos_counts},
        subtitle="→ '순함/저자극'과 '촉촉함/보습'이 압도적 - 토너 핵심 가치")

    # 슬라이드 12: Pain Points TOP 10 (동적)
    top_pains = pain_counts.most_common(8)
    pain_rows = []
    for i, (pain, cnt) in enumerate(top_pains, 1):
        pain_rows.append([str(i), pain, str(cnt), ""])

    add_table_slide(prs, "Pain Points TOP 10 (GPT 분석)",
        ["순위", "Pain Point", "건수", "주요 브랜드"],
        pain_rows,
        subtitle="→ 독도토너 Pain Point: 건조 47건, 효과부족 15건, 자극 14건")

    # 슬라이드 12-1: Pain Points 차트
    top_pain_names = [p[0] for p in top_pains[:6]]
    top_pain_counts = [p[1] for p in top_pains[:6]]
    add_bar_chart_slide(prs, "Pain Points 분포 (상위 6개)",
        categories=top_pain_names,
        series_data={'건수': top_pain_counts},
        subtitle="→ '건조/보습부족'이 압도적 1위 - 전 브랜드 공통 과제")

    # ========================================
    # Part 2. 세그먼트/페르소나 (7장)
    # ========================================
    print("[Part 2] 세그먼트 분석...")

    # 슬라이드 13: 섹션 타이틀
    add_section_slide(prs, "", "고객 세그먼트 분석")

    # 슬라이드 14: 페르소나 분류 기준
    add_content_slide(prs, "페르소나별 확정 분류 기준", [
        "▶ 1. 계획적 루틴 마스터 (Bulk Buyer)",
        "   구매 수량 ≥ 1.4개 AND 최근 150일 이내 구매",
        "",
        "▶ 2. 브랜드 충성 팬덤 (Brand Fanatic)",
        "   브랜드 집중도 ≥ 70% AND 총 구매액 ≥ 5만원",
        "",
        "▶ 3. 영리한 세일 사냥꾼 (Discount Seeker)",
        "   평균 할인율 ≥ 20% (세일 기간 집중 구매)",
        "",
        "▶ 4. 검증된 베스트셀러족 (Trend Follower)",
        "   구매 제품 평점 평균 ≥ 4.5점",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "제외 대상: 단일 품목 20개 이상 대량 구매자, 소모품만 구매한 유저",
    ], highlight_indices=[0, 3, 6, 9])

    # 슬라이드 15: 페르소나별 특성
    add_table_slide(prs, "페르소나별 RFM 특성",
        ["Persona", "Recency", "Frequency", "Monetary", "할인율", "수량"],
        [
            ["계획적 루틴 마스터", "71일", "2.9회", "229,484원", "15%", "1.66개"],
            ["브랜드 충성 팬덤", "262일", "2.43회", "167,131원", "14%", "1.38개"],
            ["검증된 베스트셀러족", "339일", "1.14회", "30,575원", "11%", "1.11개"],
            ["영리한 세일 사냥꾼", "363일", "1.11회", "28,139원", "27%", "1.13개"],
        ])

    # 슬라이드 16: 페르소나별 상관관계 인사이트
    add_content_slide(prs, "페르소나별 상관관계 핵심 인사이트", [
        "▶ 1. 계획적 루틴 마스터 (Bulk Buyer)",
        "   - 구매 수량 ↔ Recency: 35% 상관",
        "   - 한 번에 많이 사면 다음 구매까지 오래 걸림",
        "   → 마케팅: 소진 예상 시점에 맞춘 리마인드",
        "",
        "▶ 2. 영리한 세일 사냥꾼 (Discount Seeker)",
        "   - 할인율 ↔ 구매 빈도: 48% 상관",
        "   - 할인 혜택이 클수록 구매 횟수 급격히 증가",
        "   → 마케팅: 빅세일/쿠폰팩 중심 리타겟팅",
        "",
        "▶ 3. 브랜드 충성 팬덤 (Brand Fanatic)",
        "   - 스크롤 깊이 ↔ 구매 금액: 31% 상관",
        "   - 고관여 탐색이 고단가 제품 결제로 연결",
        "   → 마케팅: 브랜드 스토리 강화, 프리미엄 업셀링",
    ], highlight_indices=[0, 5, 10])

    # 슬라이드 17: 구매 확률 분석
    add_table_slide(prs, "페르소나별 평균 구매 확률",
        ["Persona", "구매 확률"],
        [
            ["계획적 루틴 마스터 (Bulk Buyer)", "85.1%"],
            ["브랜드 충성 팬덤 (Brand Fanatic)", "33.5%"],
            ["검증된 베스트셀러족 (Trend Follower)", "5.7%"],
            ["영리한 세일 사냥꾼 (Discount Seeker)", "4.6%"],
        ],
        subtitle="Key: Recency가 180일 이상이면 구매 확률 20% 미만으로 급락")

    # 슬라이드 18: 구매 확률 결정 요인
    add_content_slide(prs, "변수 간 구매 확률 결정 요인", [
        "▶ 1. Recency와 구매 확률 = 매우 강한 음의 상관",
        "   - 마지막 구매 후 시간이 지날수록 구매 확률 급락",
        "   - 특히 180일 경과 시 확률이 20% 미만으로 떨어짐",
        "",
        "▶ 2. 장바구니 담기 횟수와 구매 확률 = 45% 양의 상관",
        "   - 장바구니에 담는 행위 = 가장 강력한 '행동 신호'",
        "   - 장바구니 → 구매 전환에 집중 필요",
        "",
        "▶ 3. 할인율은 '세일 사냥꾼'의 핵심 변수",
        "   - 할인율 20% 이하 시 구매 확률 -35% 감소",
        "   - 이 그룹은 정가 마케팅 효과 없음",
    ], highlight_indices=[0, 4, 8])

    # 슬라이드 19: 세그먼트 결론
    add_quote_slide(prs,
        "\"우리는 71일 주기로 돌아오는\n계획적 루틴 마스터(Bulk Buyer)들이 먹여 살리는 구조.\n\n검증된 베스트셀러족(Trend Follower)들이\n첫 구매 후 339일 동안 돌아오지 않고 버려지는 비율이 너무 높음.\n\n이들을 브랜드 충성 팬덤(Brand Fanatic)으로\n전환시키기 위한 계획이 시급.\"")

    # ========================================
    # Part 3. 원인 분석 (12장)
    # ========================================
    print("[Part 3] 원인 분석...")

    # 슬라이드 20: 섹션 타이틀
    add_section_slide(prs, "2. 원인 분석", "데이터가 말하는 고객 이탈의 진실")

    # 슬라이드 21: Activation 분석
    add_table_slide(prs, "Activation - 봤는데 왜 '괜찮다'까지 못 가는가",
        ["분석 목적", "봐야 할 데이터", "분석 포인트", "질문"],
        [
            ["제품 이해", "상세페이지 스크롤", "25/50/75% 도달", "정보를 어디까지 소비?"],
            ["의심 구간", "FAQ/성분 영역", "체류/이탈", "어디서 불안해지는가?"],
            ["장바구니 진입", "행동 로그", "view→cart 전환", "이해했는데 왜 안 담나?"],
        ],
        subtitle="핵심 병목: 스크롤은 깊은데 장바구니 안 담음 → '지금 사야 할 이유'가 없음")

    # 슬라이드 22: 저평점 Pain Point
    add_table_slide(prs, "저평점 리뷰 Pain Point 분석 (1-2점)",
        ["브랜드", "자극/트러블", "효과없음", "건조", "향/냄새", "용기/패키지"],
        [
            ["독도토너", "30.8%", "46.2%", "7.7%", "0%", "0%"],
            ["아누아", "50.0%", "47.1%", "8.8%", "2.9%", "0%"],
            ["브링그린", "41.1%", "28.8%", "5.5%", "13.7%", "6.8%"],
            ["토리든", "29.9%", "34.3%", "9.0%", "9.0%", "7.5%"],
            ["에스네이처", "31.6%", "31.6%", "5.3%", "5.3%", "15.8%"],
        ],
        subtitle="→ 독도토너 저평점의 핵심: '효과 없음' (46.2%) - 순한 건 인정하지만 뭔가 부족")

    # 슬라이드 23: 독도토너 부정 리뷰 심층
    add_content_slide(prs, "독도토너 부정 리뷰 심층 분석 (63건, 1.5%)", [
        "▶ 포지셔닝 인식",
        "   - '무난한 입문용 토너'",
        "   - '순하고 자극 없는 토너'",
        "",
        "▶ 주요 불만 사항",
        "   1. 효과 부족 (46%): '물 같다', '뭔가 부족'",
        "   2. 자극/트러블 (31%): '예상과 달리 트러블'",
        "   3. 건조함 (8%): '보습력이 약함'",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "▶ 핵심 문제",
        "   → '자극 없는 토너(Nice to have)'로만 인식",
        "   → '꼭 필요한 토너(Must have)'로의 전환 필요",
    ], highlight_indices=[5, 6, 7, 12, 13])

    # 슬라이드 24: 독도토너 유입 분석
    add_content_slide(prs, "독도 토너는 유입 상품일까?", [
        "▶ 첫 방문 전환율 비교",
        "",
        "   | 구분                          | 첫 방문자 수 | 장바구니 전환율 |",
        "   |-------------------------------|-------------|----------------|",
        "   | 전체 사이트(자사몰) 첫 방문자   | 9,995명     | 15.19%         |",
        "   | 독도 토너 첫 조회자            | 1,845명     | 14.96%         |",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "▶ 분석 결과",
        "   - 독도 토너의 첫 방문 전환율이 자사몰 평균과 거의 동일",
        "   - 자사몰 첫 방문자 중 18.5%가 독도 토너 조회 (5명 중 1명)",
        "",
        "▶ 미션",
        "   - 상세페이지 전환율은 문제 없음",
        "   - 전체 방문 중 독도 토너 도달 비중을 늘리는 것이 과제",
    ], highlight_indices=[10, 11, 15])

    # 슬라이드 25: 스크롤 구간별 도달률
    add_table_slide(prs, "스크롤 구간별 도달률 (독도토너 상세페이지)",
        ["구간", "view", "wishlist", "cart", "purchase"],
        [
            ["0~25%", "40.84", "40.14", "40.94", "39.89"],
            ["25~50%", "29.03", "29.59", "30.06", "33.88"],
            ["50~75%", "14.87", "14.29", "14.71", "14.21"],
            ["75~90%", "9.78", "12.24", "10.87", "8.20"],
            ["90~100%", "5.49", "3.74", "3.41", "3.83"],
        ],
        subtitle="구매는 25~50%에서 가장 많이 발생 → 아하 모먼트는 중반부에 있다")

    # 슬라이드 26: 전환율 분석
    add_content_slide(prs, "독도토너 전환율 벤치마크", [
        "▶ 전환율 비교",
        "",
        "   | 비교 대상                | 구매 전환율 | 행동 전환율 |",
        "   |-------------------------|-----------|-----------|",
        "   | 토너 카테고리 평균       | 8.23%     | 30.28%    |",
        "   | 상위 10% 전환 상품       | 9.72%     | 33.31%    |",
        "   | 독도토너                 | 8.21%     | 29.26%    |",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "▶ 분석",
        "   - 토너 카테고리 평균과 비교해도 경쟁력 있음",
        "   - 상위 10% 전환 상품군과의 격차도 크지 않음",
        "",
        "▶ 결론",
        "   → 짧은 체류시간 + 준수한 전환율 = 의사결정 비용이 낮은 상품",
        "   → 초반에 아하모먼트를 부여하는 전략 필요",
    ], highlight_indices=[15, 16])

    # 슬라이드 27: 상세페이지 결론
    add_quote_slide(prs,
        "현재 상세페이지에 큰 문제는 없는 것으로 판단됨\n\n문제는 상세페이지가 아니라:\n1. '지금 사야 할 이유'가 없음\n2. '다시 찾아야 할 이유'가 없음\n3. '이 제품만의 특별함'이 없음\n\n→ 제품 자체의 리포지셔닝 필요")

    # 슬라이드 28: Retention 분석
    add_table_slide(prs, "Retention - 왜 다시 돌아오지 않는가",
        ["분석 목적", "봐야 할 데이터", "분석 포인트", "질문"],
        [
            ["관계 지속성", "Recency", "재방문 여부", "아직 관계가 살아있는가?"],
            ["루틴 신호", "Frequency", "재조회/재스크롤", "다시 고민할 만큼 남았는가?"],
            ["이탈 시점", "이탈 페이지", "마지막 페이지", "어디서 완전히 떠나는가?"],
        ],
        subtitle="Retention 병목 정의: '제품은 나쁘지 않지만, 다시 찾아야 할 이유가 없다'")

    # 슬라이드 29: 페르소나별 리텐션 - Bulk Buyer
    add_content_slide(prs, "리텐션 코호트 - 계획적 루틴 마스터 (Bulk Buyer)", [
        "▶ 리텐션 특징",
        "   - 분기당 40~50% 수준의 매우 안정적인 재구매율",
        "   - 월 단위에서 보이던 빈칸들이 메워짐",
        "",
        "▶ Recency 중앙값: 71일 (약 2.3개월)",
        "   - 이들은 약 2.3개월마다 돌아옴",
        "   - 분기 코호트에서 색이 진하게 유지 = 건강한 주기",
        "",
        "▶ 인사이트",
        "   - 매달 사진 않더라도 매 분기 최소 1회 이상 확실히 구매",
        "   - 3개월 내 미구매 시 강력한 이탈 신호로 간주",
        "",
        "▶ 액션",
        "   - 특정 분기에서 색이 연해지면 → 경쟁사 이탈 or 만족도 급감",
        "   - 소진 예상 시점(60일)에 리마인드 발송",
    ], highlight_indices=[4, 9, 10])

    # 슬라이드 30: 페르소나별 리텐션 - Brand Fanatic
    add_content_slide(prs, "리텐션 코호트 - 브랜드 충성 팬덤 (Brand Fanatic)", [
        "▶ 리텐션 특징",
        "   - 유입 후 1~2분기까지의 잔존율이 타 그룹 대비 가장 높음",
        "   - 초기 브랜드 경험이 분기 단위의 루틴으로 정착",
        "",
        "▶ Recency 중앙값: 262일 (약 9개월)",
        "   - 한 번 사면 크게(16.7만 원) 사지만, 다음 구매까지 거의 9개월",
        "   - '잊혀질 만하면 돌아오는' 고객",
        "",
        "▶ 액션",
        "   - 코호트에서 이들이 사라지지 않게 하려면",
        "   - 6개월 차쯤에 브랜드 스토리나 신제품 소식으로 '심리적 끈' 유지",
    ], highlight_indices=[4, 5, 6])

    # 슬라이드 31: 이탈 분석
    add_content_slide(prs, "독도토너 이탈자 분석", [
        "▶ 이탈 후 어디로 가는가?",
        "   1. 토리든: 54.3%",
        "   2. 에스네이처: 25.7%",
        "   3. 아비브: 12.4%",
        "   4. 토니모리: 7.6%",
        "",
        "▶ 가격 변화: 평균 +803원 (더 비싼 제품으로 이동)",
        "",
        "▶ 용량 변화: 평균 -153ml (더 작은 용량으로 이동)",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "▶ 해석",
        "   - 독도토너에서 '졸업'하는 패턴",
        "   - 입문용 → 본격 케어용으로 업그레이드",
        "   - '순함'만으로는 붙잡을 수 없음",
    ], highlight_indices=[1, 6, 8, 14, 15])

    # ========================================
    # Part 4. 전략 방향 (8장)
    # ========================================
    print("[Part 4] 전략 방향...")

    # 슬라이드 32: 섹션 타이틀
    add_section_slide(prs, "3. 전략 방향", "'무자극'을 넘어 '필수 베이스'로 리포지셔닝")

    # 슬라이드 33: 시장 트렌드 변화
    add_two_column_slide(prs, "스킨케어 시장 트렌드 변화",
        "AS-IS: '성분 중심' 구매자",
        [
            "• '순한 거 좋아요'",
            "• '자극 없으면 OK'",
            "• '성분 보고 샀어요'",
            "",
            "→ 독도토너가 강했던 시장",
        ],
        "TO-BE: '즉각적 효능 중심' 구매자",
        [
            "• '효과가 있어야 해요'",
            "• '순한 건 당연하고, 플러스알파가 필요'",
            "• '바르자마자 느껴지는 거 원해요'",
            "",
            "→ 독도토너가 약한 시장",
            "",
            "━━━━━━━━━━━━━━━━━━━━━",
            "독도토너의 문제:",
            "'순함/성분' → '즉각적 효능'으로",
            "이동한 시장에서 여전히 '순함'만 강조",
        ])

    # 슬라이드 34: Usage 분석
    top_usages = usage_counts.most_common(6)
    usage_lines = ["▶ 전체 사용법 분포"]
    for i, (usage, cnt) in enumerate(top_usages, 1):
        pct = cnt / total_reviews * 100
        usage_lines.append(f"   {i}. {usage}: {pct:.1f}% ({cnt:,}건)")

    usage_lines.extend([
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "▶ 브랜드별 닦토 사용률",
    ])

    # 닦토 사용률 정렬
    dakto_sorted = sorted(brand_usage.items(), key=lambda x: x[1]['닦토'], reverse=True)
    for brand, usage in dakto_sorted[:4]:
        usage_lines.append(f"   - {brand}: {usage['닦토']:.0f}%")

    usage_lines.extend(["", "→ 독도토너 = '닦토용 가성비템'으로 확실히 포지셔닝됨"])

    add_content_slide(prs, "사용법(Usage) 분석 - 헤비유저 패턴", usage_lines,
        highlight_indices=[1, 2, len(usage_lines)-1])

    # 슬라이드 34-1: 브랜드별 사용법 차트
    add_bar_chart_slide(prs, "브랜드별 사용법 비율",
        categories=brands,
        series_data={
            '닦토(%)': [brand_usage[b]['닦토'] for b in brands],
            '레이어링(%)': [brand_usage[b]['레이어링'] for b in brands],
            '스킨팩(%)': [brand_usage[b]['스킨팩/토너팩'] for b in brands],
        },
        subtitle="→ 독도토너: 닦토 사용 최다 (38%) | 순한 성분으로 닦토 시 자극 적음")

    # 슬라이드 35: Usage 기반 타겟
    add_content_slide(prs, "사용법 기반 타겟 세분화", [
        "▶ 닦토 유저 (38%)",
        "   - 대용량 + 순함이 핵심",
        "   - 매일 아침저녁 사용",
        "   - 소진 주기: 약 2개월",
        "",
        "▶ 레이어링 유저 (12%)",
        "   - 7스킨법 등 다겹 사용",
        "   - 촉촉함/흡수력 중시",
        "   - 소진 주기: 약 1개월",
        "",
        "▶ 스킨팩 유저 (7%)",
        "   - 진정/트러블 케어 목적",
        "   - 주 2-3회 집중 케어",
        "   - 소진 주기: 약 3개월",
        "",
        "→ 용도별 맞춤 CRM 메시지 가능",
    ], highlight_indices=[0, 5, 10, 15])

    # 슬라이드 36: Core Value Pivot
    add_two_column_slide(prs, "Core Value Pivot",
        "AS-IS",
        [
            "'자극 없는 토너 (Nice to have)'",
            "",
            "• 1020 입문용 토너",
            "• 순한 거 찾는 사람용",
            "• '물 같은 토너'",
            "",
            "문제점:",
            "→ 다음 단계로 '졸업'해버림",
            "→ 재구매 이유가 없음",
        ],
        "TO-BE",
        [
            "'고기능 케어의 완벽한 바탕 (Must have)'",
            "",
            "• 피부 휴식이 필요할 때",
            "• 어떤 제품과도 잘 어울리는 베이스",
            "• 기능성 제품의 효과를 높여주는 프라이머",
            "",
            "메시지 전환:",
            "'순해요'",
            "→ '이게 있어야 다른 게 잘 먹어요'",
        ])

    # 슬라이드 37: 타겟 재설정
    add_content_slide(prs, "Target Re-segmentation", [
        "▶ 기존 타겟",
        "   - 1020 스킨케어 입문자",
        "   - 민감성 피부 고객",
        "",
        "▶ 신규 타겟 추가",
        "",
        "   1. 그루밍 입문자 (무신사 채널)",
        "      - 남성 스킨케어 입문",
        "      - '일단 순한 거부터'",
        "",
        "   2. 미니멀리스트",
        "      - 올인원 루틴 선호",
        "      - '복잡한 거 싫어'",
        "",
        "   3. 고기능 유저의 베이스템",
        "      - 레티놀/비타민C 사용자",
        "      - '자극 완충용으로 필요'",
    ], highlight_indices=[6, 10, 14])

    # 슬라이드 38: 경쟁사 대비 전략
    add_table_slide(prs, "경쟁사 분석 기반 전략",
        ["경쟁사", "강점", "약점", "독도토너 기회"],
        [
            ["토리든", "보습, 대용량", "건조 불만 1위", "순함으로 차별화"],
            ["브링그린", "진정 효과", "트러블/자극", "민감성에서 우위"],
            ["에스네이처", "안정적 품질", "인지도 낮음", "가성비로 차별화"],
            ["토니모리", "가성비 1위", "이미지 올드", "브랜드 이미지 우위"],
            ["아누아", "프리미엄 진정", "가격 불만", "가격경쟁력 우위"],
        ],
        subtitle="독도토너 포지셔닝: '토리든의 보습 + 토니모리의 가성비 + 브링그린의 순함'")

    # 슬라이드 39: 소비 사이클 재정의
    add_content_slide(prs, "소비 사이클의 재정의", [
        "스킨케어는 '생필품'보다 '내구재'의 리듬을 따른다",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "▶ 월별 리텐션 패턴",
        "   - 1개월 후: 약 5% 재구매",
        "   - 2개월 후: 약 5% 재구매",
        "   - 3개월 후: 약 5% 재구매",
        "   - 이후 완만한 유지",
        "",
        "▶ 분기별 리텐션 패턴",
        "   - 1분기 후: 약 14% 재구매",
        "   - 2분기 후: 약 14% 재구매",
        "   - 이후 안정적 유지",
        "",
        "→ 월 단위가 아닌 분기 단위로 리텐션 관리 필요",
        "→ 3개월 주기 리마인드 전략 수립",
    ], highlight_indices=[0, 15, 16])

    # ========================================
    # Part 5. 실행 전술 (12장)
    # ========================================
    print("[Part 5] 실행 전술...")

    # 슬라이드 40: 섹션 타이틀
    add_section_slide(prs, "4. 실행 전술", "데이터 기반의 맞춤형 솔루션")

    # 슬라이드 41: CRM 전략 개요
    add_table_slide(prs, "페르소나별 CRM 전략",
        ["페르소나", "핵심 액션", "메시지 예시"],
        [
            ["계획적 루틴마스터", "소진 예상 시점 리마인드", "'토너 다 쓸 때 됐죠? 대용량 재입고'"],
            ["브랜드 충성 팬덤", "신제품/브랜드 스토리", "'독도 라인 새 제품 출시'"],
            ["영리한 세일사냥꾼", "빅세일/쿠폰팩 리타겟팅", "'오늘만 30% 할인'"],
            ["검증된 베스트셀러족", "베스트셀러 랭킹 신뢰", "'3개월 연속 1위 유지 중'"],
        ])

    # 슬라이드 42: 용도별 CRM 메시지
    add_content_slide(prs, "사용법 기반 CRM 메시지", [
        "▶ 닦토 유저 (38%)",
        "   '매일 닦토하시죠? 500ml 대용량 출시!'",
        "   '닦토할 때 가장 순한 토너, 독도'",
        "",
        "▶ 레이어링 유저 (12%)",
        "   '7스킨법 하시는 분들 주목! 여러 겹 발라도 무거움 없이'",
        "   '레이어링에 최적화된 물처럼 가벼운 텍스처'",
        "",
        "▶ 스킨팩 유저 (7%)",
        "   '스킨팩 하기 좋은 날씨네요'",
        "   '화장솜에 듬뿍, 진정 스킨팩 루틴'",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "→ 리뷰에서 추출한 사용법 키워드로 공감 형성",
    ], highlight_indices=[0, 4, 8])

    # 슬라이드 43: TF-IDF 분석 결과
    add_content_slide(prs, "TF-IDF 분석: 독도토너 재구매 핵심 동력", [
        "▶ 재구매 리뷰 TOP 키워드 (동시출현)",
        "   1. 순함: 43건 (압도적 1위)",
        "   2. 촉촉함: 19건",
        "   3. 무난함: 16건",
        "   4. 흡수력좋음: 13건",
        "",
        "▶ 재구매 vs 비재구매 리프트",
        "   - 순함: 1.0 (재구매자가 압도적으로 언급)",
        "   - 배송: 0.95 (빠른 배송도 재구매 요인)",
        "   - 흡수력좋음: 0.9",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "▶ 종합 분석: 독도토너 재구매의 핵심 동력",
        "   - 강력한 기본기 (순함): 유입과 유지 모두를 잡는 정체성",
        "   - 물류의 힘 (배송): 자사몰 빠른 배송이 재구매 결정에 영향",
        "   - 재구매 트리거: '순함', '흡수력좋음', '실패 없는 무난함'",
    ], highlight_indices=[1, 7, 14, 15, 16])

    # 슬라이드 44: 5대장 토너 비교
    add_table_slide(prs, "5대장 토너 최종 정량 분석",
        ["순위", "제품", "총 리뷰", "재구매율", "재구매주기"],
        [
            ["1", "토리든 다이브인", "3,946", "48.3%", "47일"],
            ["2", "토니모리 모찌", "1,755", "43.4%", "47일"],
            ["3", "에스네이처", "2,422", "43.3%", "43일"],
            ["4", "아비브 어성초", "1,144", "40.6%", "34일"],
            ["5", "라운드랩 독도", "1,012", "38.2%", "35일"],
        ],
        subtitle="독도토너 과제: 재구매율 5위 → 3위권 진입 목표")

    # 슬라이드 45: 채널별 전략
    add_content_slide(prs, "채널별 콘텐츠 전략", [
        "▶ 올리브영",
        "   - 리뷰 키워드: '순함', '진정', '결 개선'",
        "   - 기획 세트: 클렌저 + 토너 번들",
        "   - 프로모션: '민감성 피부 추천템'",
        "",
        "▶ 무신사",
        "   - 리뷰 키워드: '가성비', '입문용', '대용량'",
        "   - 기획 세트: 그루밍 입문 세트",
        "   - 프로모션: '남자도 스킨케어'",
        "",
        "▶ 자사몰",
        "   - 장점: 빠른 배송, 독점 혜택",
        "   - 기획: 대용량 구독 서비스",
        "   - 프로모션: '리마인드 쿠폰' (소진 예상 시점)",
    ], highlight_indices=[0, 5, 10])

    # 슬라이드 46: 상세페이지 개선
    add_content_slide(prs, "상세페이지 개선 방향", [
        "▶ 현재 문제",
        "   - 25~50% 구간에서 구매 결정",
        "   - 하지만 '지금 사야 할 이유' 부재",
        "",
        "▶ 개선 방향",
        "",
        "   1. 초반(0~25%)에 아하모먼트 배치",
        "      - '레티놀 쓰시죠? 그 전에 이거 먼저'",
        "      - '7스킨법 최적화 토너'",
        "",
        "   2. 중반(25~50%)에 신뢰 요소",
        "      - '올영 베스트셀러 3년 연속'",
        "      - '리뷰 33,000건+'",
        "",
        "   3. 후반(75~100%)에 긴급성",
        "      - '오늘 주문 시 내일 도착'",
        "      - '이번 달 한정 증정품'",
    ], highlight_indices=[6, 10, 14])

    # 슬라이드 47: Activation 병목 해결
    add_content_slide(prs, "Activation 병목 해결 전략", [
        "▶ 병목: '스크롤은 깊은데 장바구니 안 담음'",
        "",
        "▶ 해결책",
        "",
        "   1. 플로팅 CTA",
        "      - 스크롤 50% 도달 시 '장바구니 담기' 플로팅",
        "",
        "   2. 소셜 프루프 강화",
        "      - '지금 127명이 보고 있어요'",
        "      - '오늘 42명이 구매했어요'",
        "",
        "   3. 긴급성 부여",
        "      - '재입고 알림 523명 대기 중'",
        "      - '한정 수량 소진 임박'",
    ], highlight_indices=[0, 4, 7, 11])

    # 슬라이드 48: Retention 병목 해결
    add_content_slide(prs, "Retention 병목 해결 전략", [
        "▶ 병목: '제품은 괜찮은데 다시 찾을 이유가 없다'",
        "",
        "▶ 해결책",
        "",
        "   1. 소진 예측 리마인드",
        "      - 구매일 + 60일 후 알림톡",
        "      - '토너 다 쓸 때 됐죠?'",
        "",
        "   2. 크로스셀링",
        "      - 독도 토너 → 독도 로션/크림 추천",
        "      - '같이 쓰면 효과 2배'",
        "",
        "   3. 로열티 프로그램",
        "      - 3회 구매 시 VIP 등급",
        "      - 독점 할인 + 신제품 얼리버드",
    ], highlight_indices=[0, 4, 8, 12])

    # 슬라이드 49: 브랜드별 강점/약점 상세
    add_content_slide(prs, "독도토너 강점 & 약점 상세", [
        "▶ TOP 3 긍정 포인트",
        "   1. 순함/저자극 (1,062건)",
        "   2. 촉촉함/보습 (508건)",
        "   3. 재구매 의사 (273건)",
        "",
        "▶ TOP 3 불만 포인트",
        "   1. 건조/보습부족 (47건)",
        "   2. 효과 부족/미흡 (15건)",
        "   3. 자극/따가움 (14건)",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "▶ 강점 활용: '순함'을 '베이스템'으로 리포지셔닝",
        "▶ 약점 보완: '효과 없음' 인식 → 용도 제안으로 해결",
    ], highlight_indices=[1, 6, 12, 13])

    # 슬라이드 50: 재무 비교
    add_table_slide(prs, "경쟁사 재무 비교",
        ["구분", "달바", "라운드랩", "토리든", "토니모리", "아비브"],
        [
            ["매출액", "3,091억", "1,961억", "1,860억", "1,800억", "846억"],
            ["영업이익", "598억", "739억", "520억", "160억", "25억"],
            ["이익률", "19%", "38%", "28%", "9%", "3%"],
            ["특징", "물량 공세", "고효율 1위", "공격적 확장", "채널 다변화", "비용 초과"],
        ],
        subtitle="라운드랩: 업계 최고 이익률 38% → 고효율 운영의 끝판왕")

    # 슬라이드 51: 이탈 회차별 분석
    add_content_slide(prs, "이탈 회차별 행동 분석", [
        "▶ 이탈 시점별 환승 목적지",
        "   - 1~2회차 이탈: 토리든(50%) > 토니모리(10%)",
        "   - 3~4회차 이탈: 토리든(60%) > 에스네이처(20%)",
        "   - 5회차+ 이탈: 아비브 비중 증가",
        "",
        "▶ 이탈 시점별 가격 변화",
        "   - 1회차 이탈: -1,963원 (더 싼 거 삼)",
        "   - 2회차 이탈: +883원",
        "   - 3회차+ 이탈: +3,500원 이상 (업그레이드)",
        "",
        "▶ 인사이트",
        "   - 초기 이탈자: 가격에 민감, 더 저렴한 대안 탐색",
        "   - 후기 이탈자: 효능에 민감, 더 좋은 제품으로 업그레이드",
        "   - 3회 구매 이후가 '진짜 충성 고객' 전환 분기점",
    ], highlight_indices=[1, 6, 7, 8, 13])

    # ========================================
    # Part 6. 기대 효과 & 마무리 (4장)
    # ========================================
    print("[Part 6] 기대 효과 & 마무리...")

    # 슬라이드 52: 섹션 타이틀
    add_section_slide(prs, "5. 기대 효과", "KPI 및 성과 관리")

    # 슬라이드 53: KPI 설정
    add_table_slide(prs, "목표 KPI",
        ["지표", "현재", "목표", "개선폭"],
        [
            ["재구매율", "38.2%", "45%", "+6.8%p"],
            ["재구매 주기 (중앙값)", "35일", "45일 유지", "-"],
            ["자사몰 장바구니 전환율", "14.96%", "18%", "+3%p"],
            ["부정 리뷰 비율", "1.5%", "1.2%", "-0.3%p"],
            ["브랜드 검색량 증가율", "-", "+15%", "-"],
        ],
        subtitle="모니터링 주기: 월간 리뷰 + 분기 전략 점검")

    # 슬라이드 54: 핵심 인사이트 요약
    add_insight_slide(prs, "핵심 인사이트 요약", [
        "① 현상: '순함'만으로는 부족한 시대 - 저평점의 46%가 '효과 없음'",
        "② 원인: '다시 찾아야 할 이유' 부재 - Pre-Retention 상태로 전이 실패",
        "③ 전략: 'Nice to have' → 'Must have' - 고기능 케어의 베이스템으로 전환",
        "④ 전술: 용도별 맞춤 CRM - 닦토 38%, 레이어링 12% 등 사용법 기반 메시지",
        "⑤ 기대: 재구매율 38% → 45% - 5위에서 3위권 진입",
    ])

    # 슬라이드 55: 마무리
    add_quote_slide(prs,
        "독도토너의 브랜드 헤리티지(Classic)를\n현대적 감각(Essential)으로 재해석\n\n━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n\n'순한 토너'에서\n'모든 스킨케어의 시작점'으로\n\n━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n\n감사합니다.")

    # ========================================
    # 저장
    # ========================================
    output_path = 'output/독도토너_통합분석_리포트.pptx'
    prs.save(output_path)

    print("\n" + "=" * 60)
    print(f"✅ 리포트 생성 완료: {output_path}")
    print(f"   - 총 {len(prs.slides)}장 슬라이드")
    print("=" * 60)


if __name__ == "__main__":
    main()
