# -*- coding: utf-8 -*-
"""
토너 리뷰 분석 리포트 생성
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
import json
import sys
from datetime import datetime

sys.stdout.reconfigure(encoding='utf-8')

# 색상 정의
NAVY = RGBColor(0x1a, 0x36, 0x5d)
LIGHT_BLUE = RGBColor(0x4a, 0x90, 0xd9)
GREEN = RGBColor(0x27, 0xae, 0x60)
RED = RGBColor(0xe7, 0x4c, 0x3c)
ORANGE = RGBColor(0xf3, 0x9c, 0x12)
GRAY = RGBColor(0x7f, 0x8c, 0x8d)
WHITE = RGBColor(0xff, 0xff, 0xff)

def add_title_slide(prs, title, subtitle):
    """타이틀 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # 배경 박스
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(10), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_section_slide(prs, title):
    """섹션 구분 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), Inches(10), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_lines, highlight_indices=None):
    """내용 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더 바
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 내용
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = NAVY
        p.space_after = Pt(8)

        if highlight_indices and i in highlight_indices:
            p.font.bold = True
            p.font.color.rgb = LIGHT_BLUE

def add_table_slide(prs, title, headers, rows):
    """테이블 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더 바
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 테이블
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.3), Inches(1.2), Inches(9.4), Inches(0.4 * num_rows)).table

    # 헤더
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # 데이터
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = NAVY
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

def add_insight_slide(prs, title, insights):
    """인사이트 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 " + title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 인사이트 박스들
    y_pos = 1.3
    for insight in insights:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xf8, 0xf9, 0xfa)
        box.line.color.rgb = LIGHT_BLUE

        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.15), Inches(9), Inches(0.7))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = insight
        p.font.size = Pt(14)
        p.font.color.rgb = NAVY

        y_pos += 1.0


def main():
    print("슬라이드 리포트 생성 중...")

    # 데이터 로드
    df = pd.read_csv('data/merged_reviews_processed.csv', encoding='utf-8-sig')
    cat = json.load(open('output/gpt_analysis_categorized.json', encoding='utf-8'))

    # sentiment 매핑
    sent_map = {r['idx']: r for r in cat}

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== 1. 타이틀 ==========
    add_title_slide(prs,
                    "토너 리뷰 감성분석 리포트",
                    f"7개 브랜드 33,908건 GPT-4o 기반 분석 | {datetime.now().strftime('%Y.%m.%d')}")

    # ========== 2. Executive Summary ==========
    add_section_slide(prs, "Executive Summary")

    # 핵심 지표
    total_reviews = len(df)
    pos_count = sum(1 for r in cat if r.get('sentiment') == 'POS')
    neg_count = sum(1 for r in cat if r.get('sentiment') == 'NEG')
    pos_rate = pos_count / total_reviews * 100
    neg_rate = neg_count / total_reviews * 100

    add_content_slide(prs, "분석 개요", [
        f"▶ 분석 대상: 7개 토너 브랜드, 총 {total_reviews:,}건 리뷰",
        f"▶ 분석 기간: 2025.02 ~ 2026.01 (12개월)",
        f"▶ 플랫폼: 올리브영 27,745건 (82%) / 무신사 6,163건 (18%)",
        "",
        f"▶ 전체 긍정률: {pos_rate:.1f}% (부정률 {neg_rate:.1f}%)",
        f"▶ 긍정률 1위: 토니모리 (93.5%)",
        f"▶ 가장 많은 Pain Point: 건조/보습부족 (463건)",
        f"▶ 가장 많은 Positive Point: 순함/저자극 (5,495건)",
    ], highlight_indices=[4, 5, 6, 7])

    # ========== 3. 브랜드별 포지셔닝 ==========
    add_section_slide(prs, "브랜드별 포지셔닝 분석")

    # 브랜드별 감성 비교
    brand_data = []
    for brand in ['토리든', '브링그린', '독도토너', '에스네이처', '아누아', '토니모리', '아비브']:
        brand_cat = [r for r in cat if r.get('brand') == brand]
        total = len(brand_cat)
        pos = sum(1 for r in brand_cat if r.get('sentiment') == 'POS')
        neg = sum(1 for r in brand_cat if r.get('sentiment') == 'NEG')
        brand_data.append([brand, f"{total:,}", f"{pos/total*100:.1f}%", f"{neg/total*100:.1f}%"])

    add_table_slide(prs, "브랜드별 감성 분포",
                    ["브랜드", "리뷰 수", "긍정률", "부정률"],
                    brand_data)

    # 브랜드 포지셔닝 맵
    add_content_slide(prs, "브랜드 포지셔닝 맵", [
        "[ 보습 중심 ]                              [ 진정 중심 ]",
        "",
        "  토니모리 (보습 80.7%)          브링그린 (진정 59.7%)",
        "  에스네이처 (보습 73.4%)        아비브 (진정 61.3%)",
        "  토리든 (보습 70.3%)            아누아 (진정 57.3%)",
        "",
        "                    독도토너 (보습 54.9% / 진정 52.5%)",
        "                         ↑ 균형형 포지셔닝",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "[ 가성비 리더 ]                          [ 프리미엄 이미지 ]",
        "  토니모리 (가성비 39.3%)        아누아 (인생템 60.3%)",
        "  브링그린 (가성비 21.8%)        아비브 (인생템 59.3%)",
    ])

    # ========== 4. 브랜드별 강점/약점 ==========
    add_section_slide(prs, "브랜드별 강점 & 약점")

    # 브랜드별 TOP Positive/Pain (재분류 후)
    brand_analysis = {
        '토리든': {'pos': ['촉촉함/보습 (1,880)', '순함/저자극 (1,808)', '만족 (633)'],
                  'pain': ['건조/보습부족 (163)', '트러블/뾰루지 (33)', '효과 부족/미흡 (29)']},
        '브링그린': {'pos': ['순함/저자극 (686)', '진정 효과 (644)', '재구매 의사 (474)'],
                    'pain': ['건조/보습부족 (102)', '자극/따가움 (93)', '트러블/뾰루지 (43)']},
        '독도토너': {'pos': ['순함/저자극 (1,062)', '촉촉함/보습 (508)', '재구매 의사 (273)'],
                    'pain': ['건조/보습부족 (47)', '효과 부족/미흡 (15)', '자극/따가움 (14)']},
        '에스네이처': {'pos': ['촉촉함/보습 (763)', '순함/저자극 (572)', '만족 (244)'],
                      'pain': ['건조/보습부족 (31)', '눈 자극 (15)', '트러블/뾰루지 (13)']},
        '아누아': {'pos': ['순함/저자극 (681)', '촉촉함/보습 (407)', '진정 효과 (254)'],
                  'pain': ['건조/보습부족 (61)', '가격 불만 (35)', '자극/따가움 (28)']},
        '토니모리': {'pos': ['촉촉함/보습 (695)', '가성비/저렴 (480)', '대용량 (382)'],
                    'pain': ['건조/보습부족 (26)', '냄새/향 불만 (13)', '끈적임/흡수불량 (11)']},
        '아비브': {'pos': ['순함/저자극 (406)', '촉촉함/보습 (350)', '진정 효과 (272)'],
                  'pain': ['건조/보습부족 (33)', '효과 부족/미흡 (19)', '냄새/향 불만 (18)']},
    }

    for brand in ['토리든', '브링그린', '독도토너']:
        data = brand_analysis[brand]
        add_content_slide(prs, f"{brand} 강점 & 약점", [
            "◆ TOP 3 긍정 포인트",
            f"   1. {data['pos'][0]}",
            f"   2. {data['pos'][1]}",
            f"   3. {data['pos'][2]}",
            "",
            "◆ TOP 3 불만 포인트",
            f"   1. {data['pain'][0]}",
            f"   2. {data['pain'][1]}",
            f"   3. {data['pain'][2]}",
        ], highlight_indices=[1, 6])

    for brand in ['에스네이처', '아누아', '토니모리', '아비브']:
        data = brand_analysis[brand]
        add_content_slide(prs, f"{brand} 강점 & 약점", [
            "◆ TOP 3 긍정 포인트",
            f"   1. {data['pos'][0]}",
            f"   2. {data['pos'][1]}",
            f"   3. {data['pos'][2]}",
            "",
            "◆ TOP 3 불만 포인트",
            f"   1. {data['pain'][0]}",
            f"   2. {data['pain'][1]}",
            f"   3. {data['pain'][2]}",
        ], highlight_indices=[1, 6])

    # ========== 5. 플랫폼 비교 ==========
    add_section_slide(prs, "플랫폼 비교 분석")

    add_content_slide(prs, "올리브영 vs 무신사", [
        "                        올리브영              무신사",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "리뷰 수                  27,745건             6,163건",
        "전체 비중                 81.8%               18.2%",
        "",
        "긍정률                    89.7%               95.5%",
        "부정률                     2.3%                0.6%",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "▶ 무신사 리뷰가 올리브영 대비 긍정률 5.8%p 높음",
        "▶ 무신사 부정률은 올리브영의 1/4 수준",
        "▶ 무신사 주요 브랜드: 토리든, 에스네이처, 아비브, 토니모리",
    ], highlight_indices=[5, 6, 10, 11])

    # ========== 6. 월별 추이 ==========
    add_section_slide(prs, "월별 감성 추이")

    monthly_data = [
        ['2025-08', '91%', '92%', '88%', '92%', '88%', '95%', '88%'],
        ['2025-09', '92%', '88%', '92%', '93%', '91%', '92%', '92%'],
        ['2025-10', '92%', '88%', '90%', '96%', '88%', '94%', '91%'],
        ['2025-11', '91%', '88%', '89%', '91%', '89%', '91%', '94%'],
        ['2025-12', '89%', '87%', '90%', '90%', '92%', '92%', '88%'],
        ['2026-01', '89%', '88%', '90%', '93%', '-', '91%', '85%'],
    ]

    add_table_slide(prs, "브랜드별 월별 긍정률 추이 (최근 6개월)",
                    ["월", "토리든", "브링그린", "독도토너", "에스네이처", "아누아", "토니모리", "아비브"],
                    monthly_data)

    add_content_slide(prs, "월별 추이 분석", [
        "▶ 전체 긍정률 추이: 91.3% (2월) → 89.6% (1월)",
        "   - 하반기로 갈수록 소폭 하락 추세 (-1.7%p)",
        "",
        "▶ 안정적 브랜드 (변동폭 5%p 이내)",
        "   - 브링그린: 87~88% 유지",
        "   - 토니모리: 91~95% 유지",
        "   - 독도토너: 89~92% 유지",
        "",
        "▶ 변동 브랜드",
        "   - 아비브: 85~94% (9%p 변동)",
        "   - 에스네이처: 91~96% (5%p 변동)",
        "",
        "▶ 토리든: 92% → 89% (3%p 하락 추세 모니터링 필요)",
    ], highlight_indices=[0, 9, 12])

    # ========== 7. 제형/사용법 분석 ==========
    add_section_slide(prs, "제형 & 사용법 분석")

    add_content_slide(prs, "제형(Texture) 인식", [
        "▶ 전체 브랜드 '물같음' 제형이 압도적 (65~83%)",
        "",
        "  브랜드별 '물같음' 비율:",
        "   - 토리든: 82.7%",
        "   - 독도토너: 82.4%",
        "   - 에스네이처: 79.4%",
        "   - 아누아: 75.5%",
        "   - 브링그린: 73.9%",
        "   - 아비브: 72.7%",
        "   - 토니모리: 65.7% (쫀쫀함 17% 특징)",
        "",
        "▶ 토니모리만 '쫀쫀함' 제형 언급 17%로 차별화",
    ], highlight_indices=[0, 9, 11])

    add_content_slide(prs, "주요 사용법(Usage)", [
        "▶ 사용법 TOP 6",
        "",
        "   1. 닦토 (닦아내는 토너): 31.7% (10,761건)",
        "   2. 레이어링: 14.9% (5,062건)",
        "   3. 스킨팩/토너팩: 8.8% (3,000건)",
        "   4. 바디 사용: 0.6% (196건)",
        "   5. 미스트/스프레이: 0.5% (153건)",
        "   6. 흡토(패팅): 0.4% (148건)",
        "",
        "▶ 닦토 + 레이어링이 전체의 46.6%",
        "   → 대용량, 순함이 중요한 구매 요인",
    ], highlight_indices=[2, 3, 10])

    # 브랜드별 사용법 비교
    usage_data = [
        ['토리든', '30%', '18%', '9%', '1%'],
        ['브링그린', '33%', '13%', '12%', '1%'],
        ['독도토너', '38%', '12%', '7%', '1%'],
        ['에스네이처', '31%', '15%', '5%', '0%'],
        ['아누아', '31%', '13%', '10%', '1%'],
        ['토니모리', '31%', '18%', '10%', '1%'],
        ['아비브', '27%', '13%', '8%', '1%'],
    ]
    add_table_slide(prs, "브랜드별 사용법 비율",
                    ["브랜드", "닦토", "레이어링", "스킨팩", "바디"],
                    usage_data)

    add_content_slide(prs, "사용법 인사이트", [
        "▶ 독도토너: 닦토 사용 최다 (38%)",
        "   → 순한 성분으로 닦토 시 자극 적음 강조",
        "",
        "▶ 토리든/토니모리: 레이어링 사용 최다 (18%)",
        "   → 대용량 + 촉촉함으로 여러 겹 사용에 적합",
        "",
        "▶ 브링그린: 스킨팩 사용 최다 (12%)",
        "   → 진정 효과로 스킨팩 활용 마케팅 가능",
        "",
        "▶ 아비브: 닦토 사용 최저 (27%)",
        "   → 진정/트러블 케어 목적, 패팅 사용 선호",
    ], highlight_indices=[0, 3, 6, 9])

    # ========== 8. 브랜드별 부정 리뷰 분석 ==========
    add_section_slide(prs, "브랜드별 부정 리뷰 분석")

    # 부정률 비교
    neg_data = [
        ['토리든', '199건', '1.9%', '건조/보습부족, 트러블, 효과부족'],
        ['브링그린', '176건', '2.8%', '트러블, 자극/따가움, 건조'],
        ['독도토너', '63건', '1.5%', '건조, 자극, 효과부족'],
        ['에스네이처', '49건', '1.3%', '건조, 트러블, 자극'],
        ['아누아', '92건', '2.6%', '트러블, 자극, 건조'],
        ['토니모리', '45건', '1.5%', '트러블, 건조, 자극'],
        ['아비브', '55건', '2.0%', '트러블, 효과부족, 진정효과부족'],
    ]
    add_table_slide(prs, "브랜드별 부정 리뷰 현황",
                    ["브랜드", "부정 리뷰", "부정률", "주요 불만"],
                    neg_data)

    add_content_slide(prs, "부정 리뷰 핵심 분석", [
        "▶ 부정률 높은 브랜드",
        "   - 브링그린 (2.8%): 트러블/자극 불만 집중",
        "   - 아누아 (2.6%): 트러블 + 가격 불만 복합",
        "",
        "▶ 부정률 낮은 브랜드",
        "   - 에스네이처 (1.3%): 안정적 품질 인식",
        "   - 독도토너/토니모리 (1.5%): 기대치 충족",
        "",
        "▶ 공통 불만 패턴",
        "   - 건조/보습부족: 전 브랜드 공통 (특히 토리든)",
        "   - 트러블/뾰루지: 브링그린, 아누아에서 심각",
        "   - 자극/따가움: 민감성 피부 불만 집중",
    ], highlight_indices=[1, 2, 5, 6, 9, 10, 11])

    # 브랜드별 상세 분석
    add_content_slide(prs, "토리든 부정 리뷰 분석 (199건, 1.9%)", [
        "▶ 포지셔닝 인식",
        "   - 보습/촉촉함 기대가 높은 브랜드",
        "   - 대용량 가성비 토너로 인식",
        "",
        "▶ 주요 불만 사항",
        "   1. 건조/보습부족 (46건): '생각보다 건조', '속건조'",
        "   2. 트러블/뾰루지 (24건): '피부가 뒤집어짐'",
        "   3. 효과 부족 (11건): '기대한 효과 없음'",
        "",
        "▶ 개선 포인트",
        "   → 보습 지속력 강화, 건성 피부용 라인 검토",
    ], highlight_indices=[5, 6, 7, 10])

    add_content_slide(prs, "브링그린 부정 리뷰 분석 (176건, 2.8%)", [
        "▶ 포지셔닝 인식",
        "   - 진정/트러블 케어 전문 브랜드",
        "   - 민감성 피부용으로 인식",
        "",
        "▶ 주요 불만 사항",
        "   1. 트러블/뾰루지 (32건): '오히려 트러블 발생'",
        "   2. 자극/따가움 (24건): '민감해서 따가움'",
        "   3. 건조/보습부족 (18건): '진정은 되는데 건조'",
        "",
        "▶ 개선 포인트",
        "   → 민감성 피부 테스트 강화, 성분 재검토",
    ], highlight_indices=[5, 6, 7, 10])

    add_content_slide(prs, "아누아 부정 리뷰 분석 (92건, 2.6%)", [
        "▶ 포지셔닝 인식",
        "   - 프리미엄 진정 토너",
        "   - 어성초 진정 효과 기대",
        "",
        "▶ 주요 불만 사항",
        "   1. 트러블/뾰루지 (16건): '트러블 개선 안됨'",
        "   2. 자극/따가움 (10건): '생각보다 자극적'",
        "   3. 가격 불만 (9건): '가격 대비 효과 의문'",
        "",
        "▶ 개선 포인트",
        "   → 가격 대비 가치 커뮤니케이션 강화",
    ], highlight_indices=[5, 6, 7, 10])

    # ========== 9. 핵심 인사이트 ==========
    add_section_slide(prs, "핵심 인사이트 & 제언")

    add_insight_slide(prs, "브랜드 전략 제언", [
        "① 토리든: 점유율 1위(30%). 부정 리뷰 중 '건조' 불만 최다 → 보습력 강화 필요",
        "② 브링그린: 진정 효과 강점이나 '트러블/자극' 불만 → 민감성 테스트 강화",
        "③ 독도토너: '순함/저자극' 강점 유지. 닦토 사용률 38%로 최고",
        "④ 아누아: 트러블 케어 기대 불충족 + 가격 불만 → 효능/가치 재검토",
        "⑤ 토니모리: 가성비 리더. 부정률 1.5%로 안정적 품질 유지",
    ])

    add_insight_slide(prs, "제품 개선 제언", [
        "① 공통 Pain Point: '건조/보습부족' 전 브랜드 1위 → 보습 지속력 강화",
        "② 트러블 케어 브랜드: 역설적 트러블 발생 불만 → 성분 재검토 필요",
        "③ 사용법 마케팅: 닦토(31.7%) + 레이어링(14.9%) 중심 콘텐츠 강화",
        "④ 민감성 시장: 자극/따가움 불만 많음 → 패치 테스트 권장 안내",
        "⑤ 시즌별 전략: 겨울철 건조 불만 증가 대비 → 고보습 라인 준비",
    ])

    # ========== 10. 부록 ==========
    add_section_slide(prs, "부록: 상세 데이터")

    # 전체 Pain Points TOP 10 (재분류 후)
    pain_data = [
        ['건조/보습부족', '463', '전체 (토리든 163, 브링그린 102)'],
        ['자극/따가움', '181', '브링그린, 아누아'],
        ['트러블/뾰루지', '139', '브링그린, 토리든'],
        ['효과 부족/미흡', '115', '토리든, 아비브'],
        ['가격 불만', '101', '아누아'],
        ['눈 자극', '88', '에스네이처, 브링그린'],
        ['진정효과 부족', '67', '아누아, 아비브'],
        ['냄새/향 불만', '64', '토니모리, 아비브'],
        ['끈적임/흡수불량', '49', '토니모리'],
        ['재구매 의사 없음', '38', '토리든'],
    ]
    add_table_slide(prs, "전체 Pain Points TOP 10",
                    ["Pain Point", "건수", "주요 브랜드"],
                    pain_data)

    # 전체 Positive Points TOP 10 (재분류 후)
    pos_data = [
        ['순함/저자극', '5,495', '토리든, 독도토너, 브링그린'],
        ['촉촉함/보습', '5,066', '토리든, 에스네이처, 토니모리'],
        ['재구매 의사', '2,138', '토리든, 브링그린'],
        ['만족', '2,024', '토리든, 에스네이처'],
        ['가성비/저렴', '1,924', '토니모리, 브링그린'],
        ['대용량', '1,751', '토니모리, 토리든'],
        ['진정 효과', '1,492', '브링그린, 아비브'],
        ['산뜻함/비끈적', '1,293', '토리든, 에스네이처'],
        ['무난함', '1,181', '독도토너, 에스네이처'],
        ['추천', '803', '전체'],
    ]
    add_table_slide(prs, "전체 Positive Points TOP 10",
                    ["Positive Point", "건수", "주요 브랜드"],
                    pos_data)

    # ========== 저장 ==========
    output_path = 'output/토너_리뷰분석_리포트.pptx'
    prs.save(output_path)
    print(f"\n✅ 리포트 생성 완료: {output_path}")
    print(f"   - 총 {len(prs.slides)}장 슬라이드")


if __name__ == "__main__":
    main()
