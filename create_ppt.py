# -*- coding: utf-8 -*-
"""
분석 결과 PPT 생성
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# PPT 생성
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 부제목
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(title, content_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # 내용
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(6))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    return slide

def add_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # 테이블
    cols = len(headers)
    table = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.5 * (len(rows) + 1))).table

    # 헤더
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    # 데이터
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(13)

    return slide

def add_image_slide(title, image_path, caption=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    # 이미지
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(0.9), width=Inches(12.333))
    else:
        # 이미지 없으면 텍스트로 표시
        box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[이미지: {image_path}]"
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_two_images_slide(title, img1_path, img2_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    # 이미지 2개
    if os.path.exists(img1_path):
        slide.shapes.add_picture(img1_path, Inches(0.3), Inches(0.9), width=Inches(6.3))
    if os.path.exists(img2_path):
        slide.shapes.add_picture(img2_path, Inches(6.8), Inches(0.9), width=Inches(6.3))

    return slide

# ===== 슬라이드 생성 =====

# 1. 표지
add_title_slide(
    "올리브영 토너 리뷰 분석",
    "7개 브랜드 27,745건 리뷰 심층 분석 | v2.0"
)

# 2. 분석 개요
add_content_slide("분석 개요", [
    "분석 대상: 올리브영 토너 카테고리 7개 브랜드",
    "",
    "• 토리든: 8,065건",
    "• 브링그린: 6,312건",
    "• 독도토너: 3,750건",
    "• 아누아: 3,544건",
    "• 에스네이처: 2,232건",
    "• 토니모리: 2,195건",
    "• 아비브: 1,647건",
    "",
    "총 27,745건 리뷰 분석"
])

# 3. v2.0 개선사항
add_content_slide("v2.0 분석 로직 개선", [
    "기존 문제점: 별점 5점이지만 실제로는 부정적인 리뷰를 POS로 분류",
    "",
    "개선 내용:",
    "• 피부질병 키워드 탐지 (모낭염, 알러지 등 24개)",
    "• 역접 패턴 처리 (~했었으나, ~지만, ~는데)",
    "• 중단/사용중지 키워드 탐지",
    "• 과거 사용 패턴 필터링",
    "• 부정 문맥 내 태그 오탐 방지",
    "",
    "결과: NEG 탐지 230건 → 3,937건 (+3,707건)"
])

# 4. 감성 분포 변화
add_table_slide("감성 분포 변화 (v1 → v2)",
    ["감성", "v1 (이전)", "v2 (개선)", "변화"],
    [
        ["POS", "26,747건 (96.4%)", "23,251건 (83.8%)", "-3,496건"],
        ["NEU", "768건 (2.8%)", "557건 (2.0%)", "-211건"],
        ["NEG", "230건 (0.8%)", "3,937건 (14.2%)", "+3,707건"],
    ]
)

# 5. 문제 케이스 예시
add_content_slide("문제 케이스 검증", [
    "리뷰: \"예전에 닦토로 잘써서 재구매했었으나 모낭염이 올라와서 중단했네여요\"",
    "별점: 5점",
    "",
    "[v1 분석 결과] - 잘못됨",
    "• 감성: POS",
    "• 사용법: ['닦토']",
    "• 가치: ['인생템']",
    "",
    "[v2 분석 결과] - 정확함",
    "• 감성: NEG",
    "• 사용법: []",
    "• 가치: []",
])

# 6. 브랜드별 NEG 비율
add_table_slide("브랜드별 NEG 비율 (v2.0)",
    ["브랜드", "리뷰 수", "평균 평점", "NEG 비율", "피부질병 언급률"],
    [
        ["아누아", "3,544", "4.8", "21.3%", "13.0%"],
        ["브링그린", "6,312", "4.8", "20.7%", "12.6%"],
        ["아비브", "1,647", "4.8", "16.2%", "9.0%"],
        ["독도토너", "3,750", "4.9", "11.1%", "3.3%"],
        ["토리든", "8,065", "4.8", "10.3%", "3.7%"],
        ["에스네이처", "2,232", "4.8", "9.6%", "2.8%"],
        ["토니모리", "2,195", "4.9", "6.9%", "2.2%"],
    ]
)

# 7. 종합 대시보드
add_image_slide("종합 대시보드", "output/figures/dashboard_summary.png")

# 8. Pain Point 분석
add_two_images_slide("Pain Point 분석",
    "output/figures/pain_point_heatmap.png",
    "output/figures/pain_point_by_brand.png"
)

# 9. 재구매/이탈 분석
add_two_images_slide("재구매/이탈 분석",
    "output/figures/loyalty_comparison.png",
    "output/figures/net_loyalty.png"
)

# 10. 포지셔닝 분석
add_two_images_slide("포지셔닝 분석",
    "output/figures/3B_positioning_benefit.png",
    "output/figures/3B_positioning_texture.png"
)

# 11. 브랜드 포지셔닝 레이더
add_image_slide("브랜드 포지셔닝", "output/figures/brand_positioning_radar.png")

# 12. 월별 트렌드
add_two_images_slide("월별 키워드 트렌드",
    "output/figures/monthly_trend_all.png",
    "output/figures/monthly_trend_by_brand.png"
)

# 13. 피부타입 분포
add_image_slide("피부타입 분포", "output/figures/skin_type_distribution.png")

# 14. 전환 매트릭스
add_image_slide("고객 전환 매트릭스", "output/figures/3E_switch_matrix.png")

# 15. 핵심 인사이트
add_content_slide("핵심 인사이트", [
    "1. 진정 브랜드의 딜레마",
    "   - 아누아, 브링그린: 진정 특화 but 피부질병 언급률 10% 이상",
    "   - 민감성 타겟 → 자극 리포트 다수",
    "",
    "2. 별점의 한계",
    "   - 평균 별점 4.8~4.9로 차이 미미",
    "   - 실제 NEG 비율은 6.9%~21.3%로 3배 이상 차이",
    "",
    "3. 토니모리의 안정성",
    "   - 가장 낮은 NEG(6.9%), 피부질병(2.2%)",
    "   - 가성비 + 낮은 자극으로 대중적 안전성 확보"
])

# 16. 브랜드별 제언
add_table_slide("브랜드별 액션 아이템",
    ["브랜드", "현황", "제언"],
    [
        ["아누아", "NEG 21.3%", "포뮬러 순화, 패치테스트 권고"],
        ["브링그린", "NEG 20.7%", "어성초 농도 조절, 민감성 라인 분리"],
        ["아비브", "NEG 16.2%", "진정 효과 홍보 강화"],
        ["토리든", "NEG 10.3%", "현상 유지, 보습 마케팅 강화"],
        ["토니모리", "NEG 6.9%", "가성비 포지션 유지"],
    ]
)

# 17. 마무리
add_title_slide(
    "감사합니다",
    "분석 대상: 7개 브랜드 27,745건 리뷰"
)

# 저장
output_path = "docs/review_analysis_v2.pptx"
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
